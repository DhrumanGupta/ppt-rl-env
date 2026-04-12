import pytest
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches

from ppt_agent.server.utils.pptx_functions import PptxEditor


@pytest.fixture
def editor():
    return PptxEditor()


def add_blank_slide(editor: PptxEditor) -> int:
    return editor.add_slide()


def add_textbox(editor: PptxEditor, slide_index: int, text: str = "sample") -> int:
    slide = editor.prs.slides[slide_index]
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text_frame.text = text
    return len(slide.shapes) - 1


def test_add_slide_creates_slide(editor):
    slide_index = editor.add_slide()

    assert slide_index == 0
    assert len(editor.prs.slides) == 1


def test_validate_slide_raises_for_invalid_indexes(editor):
    add_blank_slide(editor)

    with pytest.raises(IndexError, match="Slide index -1 out of range"):
        editor._validate_slide(-1)

    with pytest.raises(IndexError, match="Slide index 1 out of range"):
        editor._validate_slide(1)


def test_get_shape_raises_for_invalid_shape_index(editor):
    slide_index = add_blank_slide(editor)

    with pytest.raises(IndexError, match="Shape index 0 out of range"):
        editor._get_shape(slide_index, 0)


def test_reorder_slide_moves_slide_to_new_position(editor):
    for _ in range(3):
        add_blank_slide(editor)

    original_ids = [slide.slide_id for slide in editor.prs.slides]

    editor.reorder_slide(0, 2)

    reordered_ids = [slide.slide_id for slide in editor.prs.slides]
    assert reordered_ids == [original_ids[1], original_ids[2], original_ids[0]]


def test_reorder_slide_raises_for_invalid_new_index(editor):
    add_blank_slide(editor)

    with pytest.raises(IndexError, match="new_index out of bounds"):
        editor.reorder_slide(0, 1)


def test_insert_text_replaces_existing_text_in_textbox(editor):
    slide_index = add_blank_slide(editor)
    shape_index = add_textbox(editor, slide_index, "Before")

    editor.insert_text(slide_index, shape_index, "After")

    assert editor.prs.slides[slide_index].shapes[shape_index].text_frame.text == "After"


def test_insert_text_preserves_newline_separated_paragraphs(editor):
    slide_index = add_blank_slide(editor)
    shape_index = add_textbox(editor, slide_index)

    editor.insert_text(slide_index, shape_index, "One\nTwo\nThree")

    paragraphs = (
        editor.prs.slides[slide_index].shapes[shape_index].text_frame.paragraphs
    )
    assert [paragraph.text for paragraph in paragraphs] == ["One", "Two", "Three"]


def test_insert_text_rejects_non_text_shape(editor):
    slide_index = add_blank_slide(editor)
    slide = editor.prs.slides[slide_index]
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(2)
    )

    with pytest.raises(ValueError, match="Shape has no text frame"):
        editor.insert_text(slide_index, 0, "One")


def test_add_chart_uses_default_chart_type_for_unknown_values(editor):
    slide_index = add_blank_slide(editor)
    slide = editor.prs.slides[slide_index]

    editor.add_chart(
        slide_index,
        "not_a_chart_type",
        {"categories": ["A", "B"], "series": [{"name": "Series 1", "values": [1, 2]}]},
        1,
        1,
        4,
        3,
    )

    assert len(slide.shapes) == 1
    assert slide.shapes[0].has_chart
    assert slide.shapes[0].chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED


def test_add_chart_applies_optional_style(editor):
    slide_index = add_blank_slide(editor)
    slide = editor.prs.slides[slide_index]

    editor.add_chart(
        slide_index,
        "column_clustered",
        {"categories": ["A", "B"], "series": [{"name": "Series 1", "values": [1, 2]}]},
        1,
        1,
        4,
        3,
        style={
            "title": "Styled Chart",
            "title_font_name": "Arial",
            "title_font_size_pt": 20,
            "title_color_hex": "#112233",
            "series_colors": ["#445566"],
        },
    )

    chart = slide.shapes[0].chart
    assert chart.chart_title.text_frame.text == "Styled Chart"
    assert chart.chart_title.text_frame.paragraphs[0].font.name == "Arial"
    assert chart.chart_title.text_frame.paragraphs[0].font.size.pt == 20
    assert str(chart.chart_title.text_frame.paragraphs[0].font.color.rgb) == "112233"
    assert str(chart.series[0].format.fill.fore_color.rgb) == "445566"


def test_style_chart_updates_existing_chart(editor):
    slide_index = add_blank_slide(editor)
    slide = editor.prs.slides[slide_index]
    editor.add_chart(
        slide_index,
        "column_clustered",
        {"categories": ["A", "B"], "series": [{"name": "Series 1", "values": [1, 2]}]},
        1,
        1,
        4,
        3,
    )

    editor.style_chart(
        slide_index,
        0,
        title="Updated Title",
        title_font_name="Calibri",
        title_font_size_pt=18,
        title_bold=True,
        title_color_hex="#AA5500",
        series_colors=["#003366"],
    )

    chart = slide.shapes[0].chart
    title_font = chart.chart_title.text_frame.paragraphs[0].font
    assert chart.chart_title.text_frame.text == "Updated Title"
    assert title_font.name == "Calibri"
    assert title_font.size.pt == 18
    assert title_font.bold is True
    assert str(title_font.color.rgb) == "AA5500"
    assert str(chart.series[0].format.fill.fore_color.rgb) == "003366"


def test_add_table_populates_available_cells_only(editor):
    slide_index = add_blank_slide(editor)
    slide = editor.prs.slides[slide_index]

    editor.add_table(
        slide_index,
        2,
        2,
        [["A1", "A2", "A3"], ["B1"]],
        1,
        1,
        4,
        2,
    )

    table = slide.shapes[0].table
    assert table.cell(0, 0).text == "A1"
    assert table.cell(0, 1).text == "A2"
    assert table.cell(1, 0).text == "B1"
    assert table.cell(1, 1).text == ""


def test_add_table_applies_optional_style(editor):
    slide_index = add_blank_slide(editor)
    slide = editor.prs.slides[slide_index]

    editor.add_table(
        slide_index,
        2,
        2,
        [["Header 1", "Header 2"], ["Value 1", "Value 2"]],
        1,
        1,
        4,
        2,
        style={
            "header_fill_hex": "#112233",
            "body_fill_hex": "#F1F5F9",
            "header_font_name": "Arial",
            "body_font_name": "Calibri",
            "header_font_size_pt": 18,
            "body_font_size_pt": 14,
            "header_font_color_hex": "#FFFFFF",
            "body_font_color_hex": "#334455",
        },
    )

    table = slide.shapes[0].table
    assert str(table.cell(0, 0).fill.fore_color.rgb) == "112233"
    assert str(table.cell(1, 0).fill.fore_color.rgb) == "F1F5F9"
    assert table.cell(0, 0).text_frame.paragraphs[0].font.name == "Arial"
    assert table.cell(1, 0).text_frame.paragraphs[0].font.name == "Calibri"
    assert table.cell(0, 0).text_frame.paragraphs[0].font.size.pt == 18
    assert table.cell(1, 0).text_frame.paragraphs[0].font.size.pt == 14
    assert str(table.cell(0, 0).text_frame.paragraphs[0].font.color.rgb) == "FFFFFF"
    assert str(table.cell(1, 0).text_frame.paragraphs[0].font.color.rgb) == "334455"


def test_style_table_updates_existing_table(editor):
    slide_index = add_blank_slide(editor)
    slide = editor.prs.slides[slide_index]
    editor.add_table(
        slide_index,
        2,
        2,
        [["Header 1", "Header 2"], ["Value 1", "Value 2"]],
        1,
        1,
        4,
        2,
    )

    editor.style_table(
        slide_index,
        0,
        header_fill_hex="#0F172A",
        body_fill_hex="#E2E8F0",
        header_font_name="Georgia",
        body_font_name="Verdana",
        header_font_size_pt=16,
        body_font_size_pt=12,
        header_font_color_hex="#FAFAFA",
        body_font_color_hex="#1E293B",
    )

    table = slide.shapes[0].table
    assert str(table.cell(0, 1).fill.fore_color.rgb) == "0F172A"
    assert str(table.cell(1, 1).fill.fore_color.rgb) == "E2E8F0"
    assert table.cell(0, 1).text_frame.paragraphs[0].font.name == "Georgia"
    assert table.cell(1, 1).text_frame.paragraphs[0].font.name == "Verdana"
    assert str(table.cell(0, 1).text_frame.paragraphs[0].font.color.rgb) == "FAFAFA"
    assert str(table.cell(1, 1).text_frame.paragraphs[0].font.color.rgb) == "1E293B"


def test_style_update_applies_font_settings_to_existing_runs(editor):
    slide_index = add_blank_slide(editor)
    shape_index = add_textbox(editor, slide_index, "")
    paragraph = (
        editor.prs.slides[slide_index].shapes[shape_index].text_frame.paragraphs[0]
    )
    run = paragraph.add_run()
    run.text = "Styled"

    editor.style_update(
        slide_index,
        shape_index,
        font_name="Arial",
        font_size_pt=18,
        bold=True,
        italic=True,
        color_hex="#112233",
    )

    assert run.font.name == "Arial"
    assert run.font.size.pt == 18
    assert run.font.bold is True
    assert run.font.italic is True
    assert str(run.font.color.rgb) == "112233"


def test_add_textbox_enables_word_wrap_by_default(editor):
    slide_index = add_blank_slide(editor)
    slide_id = editor.get_slide_id(slide_index)
    shape_id = editor.add_textbox_by_id(slide_id, 1, 1, 4, 1)
    editor.insert_text_by_id(slide_id, shape_id, "Wrapped by default")

    text_frame = editor._get_shape_by_id(slide_id, shape_id).text_frame

    assert text_frame.word_wrap is True


def test_style_text_applies_wrap_and_paragraph_spacing(editor):
    slide_index = add_blank_slide(editor)
    shape_index = add_textbox(editor, slide_index, "Line 1\nLine 2")

    editor.style_update(
        slide_index,
        shape_index,
        font_name="Arial",
        font_size_pt=18,
        color_hex="#112233",
        word_wrap=False,
        space_before_pt=3,
        space_after_pt=6,
        line_spacing=1.4,
    )

    paragraphs = (
        editor.prs.slides[slide_index].shapes[shape_index].text_frame.paragraphs
    )
    text_frame = editor.prs.slides[slide_index].shapes[shape_index].text_frame

    assert text_frame.word_wrap is False
    assert paragraphs[0].space_before.pt == 3
    assert paragraphs[0].space_after.pt == 6
    assert paragraphs[0].line_spacing == 1.4


def test_style_update_rejects_non_text_shape(editor):
    slide_index = add_blank_slide(editor)
    slide = editor.prs.slides[slide_index]
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(2)
    )

    with pytest.raises(ValueError, match="Shape has no text frame"):
        editor.style_update(slide_index, 0, bold=True)


def test_style_table_rejects_non_table_shape(editor):
    slide_index = add_blank_slide(editor)
    shape_index = add_textbox(editor, slide_index, "Not a table")

    with pytest.raises(ValueError, match="Shape has no table"):
        editor.style_table(slide_index, shape_index, header_fill_hex="#000000")


def test_style_chart_rejects_non_chart_shape(editor):
    slide_index = add_blank_slide(editor)
    shape_index = add_textbox(editor, slide_index, "Not a chart")

    with pytest.raises(ValueError, match="Shape has no chart"):
        editor.style_chart(slide_index, shape_index, title="Nope")


def test_add_image_raises_for_missing_file(editor):
    slide_index = add_blank_slide(editor)

    with pytest.raises(FileNotFoundError, match="Image path invalid"):
        editor.add_image(slide_index, "/tmp/does-not-exist.png", 1, 1)
