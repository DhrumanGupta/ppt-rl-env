import sys
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from src.tools.pptx_tools import PptxEditor


OUTPUT_DIR = ROOT / "outputs" / "test_slides"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def apply_slide_background(slide, color_hex: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)


def add_accent_bar(slide, color_hex: str):
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(0.35),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(color_hex)
    accent.line.fill.background()


def add_textbox(
    editor: PptxEditor, slide_index: int, x: float, y: float, w: float, h: float
) -> int:
    slide = editor.prs.slides[slide_index]
    slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return len(slide.shapes) - 1


def add_text_block(
    editor: PptxEditor,
    slide_index: int,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_name: str,
    font_size: int,
    color_hex: str,
    bold: bool = False,
    italic: bool = False,
) -> int:
    shape_index = add_textbox(editor, slide_index, x, y, w, h)
    editor.insert_text(slide_index, shape_index, text)
    editor.style_update(
        slide_index,
        shape_index,
        font_name=font_name,
        font_size_pt=font_size,
        bold=bold,
        italic=italic,
        color_hex=color_hex,
    )
    return shape_index


def build_modern_business_deck(output_path: Path):
    editor = PptxEditor()
    theme = {
        "bg": "#F6F9FC",
        "accent": "#0F62FE",
        "primary": "#102A43",
        "secondary": "#486581",
        "font": "Aptos",
    }

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, theme["bg"])
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Northstar Growth Plan 2026",
        0.9,
        0.8,
        10.5,
        0.8,
        font_name=theme["font"],
        font_size=28,
        color_hex=theme["primary"],
        bold=True,
    )
    add_text_block(
        editor,
        slide_index,
        "A focused strategy for revenue expansion, retention, and operational leverage.",
        0.9,
        1.8,
        10.5,
        0.8,
        font_name=theme["font"],
        font_size=16,
        color_hex=theme["secondary"],
    )

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, "#FFFFFF")
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Executive Summary",
        0.8,
        0.7,
        4.5,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    add_text_block(
        editor,
        slide_index,
        "Expand into two adjacent verticals\nRaise enterprise retention from 88% to 93%\nReduce onboarding time by 35% with automation",
        0.9,
        1.6,
        6.0,
        2.8,
        font_name=theme["font"],
        font_size=18,
        color_hex=theme["secondary"],
    )
    editor.add_citation(slide_index, "Source: Internal planning memo, Q1 2026")

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, "#FFFFFF")
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Quarterly Revenue Outlook",
        0.8,
        0.7,
        5.0,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    editor.add_chart(
        slide_index,
        "column_clustered",
        {
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": [
                {"name": "Baseline", "values": [18, 21, 23, 25]},
                {"name": "Target", "values": [19, 24, 28, 32]},
            ],
        },
        0.9,
        1.6,
        7.5,
        4.1,
        style={
            "title": "Revenue trajectory ($M)",
            "title_font_name": theme["font"],
            "title_font_size_pt": 18,
            "title_bold": True,
            "title_color_hex": theme["primary"],
            "legend_font_name": theme["font"],
            "legend_font_size_pt": 11,
            "legend_color_hex": theme["secondary"],
            "axis_font_name": theme["font"],
            "axis_font_size_pt": 11,
            "axis_color_hex": theme["secondary"],
            "series_colors": [theme["secondary"], theme["accent"]],
        },
    )
    add_text_block(
        editor,
        slide_index,
        "Target upside comes from enterprise upsell and partner-sourced pipeline.",
        8.7,
        2.0,
        3.2,
        1.8,
        font_name=theme["font"],
        font_size=16,
        color_hex=theme["secondary"],
    )

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, "#FFFFFF")
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Key Initiatives",
        0.8,
        0.7,
        4.5,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    editor.add_table(
        slide_index,
        4,
        3,
        [
            ["Initiative", "Owner", "Impact"],
            ["Pricing refresh", "Revenue Ops", "+6% ARR"],
            ["Guided onboarding", "Product", "-35% setup time"],
            ["Partner motion", "Sales", "+$4M pipeline"],
        ],
        0.9,
        1.6,
        8.5,
        2.5,
        style={
            "header_fill_hex": theme["accent"],
            "body_fill_hex": "#EAF2FF",
            "header_font_name": theme["font"],
            "body_font_name": theme["font"],
            "header_font_size_pt": 14,
            "body_font_size_pt": 12,
            "header_font_color_hex": "#FFFFFF",
            "body_font_color_hex": theme["primary"],
        },
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    editor.prs.save(output_path)


def build_dark_analytics_deck(output_path: Path):
    editor = PptxEditor()
    theme = {
        "bg": "#0B1020",
        "surface": "#131A2A",
        "accent": "#4FD1C5",
        "primary": "#F7FAFC",
        "secondary": "#A0AEC0",
        "font": "Aptos Display",
    }

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, theme["bg"])
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Platform Reliability Deep Dive",
        0.8,
        0.9,
        8.5,
        0.8,
        font_name=theme["font"],
        font_size=30,
        color_hex=theme["primary"],
        bold=True,
    )
    add_text_block(
        editor,
        slide_index,
        "Telemetry review of latency, incident rate, and regional capacity.",
        0.8,
        1.9,
        9.0,
        0.7,
        font_name=theme["font"],
        font_size=16,
        color_hex=theme["secondary"],
    )

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, theme["surface"])
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "System Status Snapshot",
        0.8,
        0.7,
        5.0,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    add_text_block(
        editor,
        slide_index,
        "Median latency improved by 18%\nError budget burn stabilized after week 6\nEU capacity remains the main scaling constraint",
        0.9,
        1.6,
        5.8,
        2.6,
        font_name=theme["font"],
        font_size=18,
        color_hex=theme["secondary"],
    )

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, theme["surface"])
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Incident Trend by Month",
        0.8,
        0.7,
        5.0,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    editor.add_chart(
        slide_index,
        "line_markers",
        {
            "categories": ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            "series": [
                {"name": "P1 incidents", "values": [9, 8, 7, 5, 4, 3]},
                {"name": "P2 incidents", "values": [18, 17, 15, 13, 11, 10]},
            ],
        },
        0.9,
        1.6,
        7.2,
        4.2,
        style={
            "title": "Monthly incident volume",
            "title_font_name": theme["font"],
            "title_font_size_pt": 18,
            "title_bold": True,
            "title_color_hex": theme["primary"],
            "legend_font_name": theme["font"],
            "legend_font_size_pt": 11,
            "legend_color_hex": theme["secondary"],
            "axis_font_name": theme["font"],
            "axis_font_size_pt": 11,
            "axis_color_hex": theme["secondary"],
            "series_colors": [theme["accent"], "#F6AD55"],
        },
    )
    add_text_block(
        editor,
        slide_index,
        "The reliability program is reducing both severity and recurrence.",
        8.5,
        2.1,
        3.2,
        1.8,
        font_name=theme["font"],
        font_size=16,
        color_hex=theme["secondary"],
    )

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, theme["surface"])
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Regional Capacity Allocation",
        0.8,
        0.7,
        5.3,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    editor.add_table(
        slide_index,
        4,
        3,
        [
            ["Region", "Utilization", "Headroom"],
            ["US-East", "68%", "High"],
            ["EU-West", "87%", "Low"],
            ["APAC", "59%", "Medium"],
        ],
        0.9,
        1.6,
        8.0,
        2.5,
        style={
            "header_fill_hex": theme["accent"],
            "body_fill_hex": "#1B2437",
            "header_font_name": theme["font"],
            "body_font_name": theme["font"],
            "header_font_size_pt": 14,
            "body_font_size_pt": 12,
            "header_font_color_hex": theme["bg"],
            "body_font_color_hex": theme["primary"],
        },
    )
    editor.add_citation(slide_index, "Source: Weekly reliability metrics")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    editor.prs.save(output_path)


def build_academic_report_deck(output_path: Path):
    editor = PptxEditor()
    theme = {
        "bg": "#F8F5EF",
        "accent": "#7A5C3E",
        "primary": "#2D241D",
        "secondary": "#6B5B4D",
        "font": "Georgia",
    }

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, theme["bg"])
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Urban Tree Canopy Study",
        0.9,
        0.9,
        8.0,
        0.8,
        font_name=theme["font"],
        font_size=28,
        color_hex=theme["primary"],
        bold=True,
    )
    add_text_block(
        editor,
        slide_index,
        "Findings from a five-year survey of neighborhood heat resilience.",
        0.9,
        1.9,
        8.8,
        0.7,
        font_name=theme["font"],
        font_size=16,
        color_hex=theme["secondary"],
        italic=True,
    )

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, "#FFFDF9")
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Research Questions",
        0.8,
        0.7,
        4.5,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    add_text_block(
        editor,
        slide_index,
        "How does canopy density affect surface temperature?\nWhich districts show the largest seasonal benefit?\nWhat investments provide the fastest public-health return?",
        0.9,
        1.6,
        6.3,
        2.8,
        font_name=theme["font"],
        font_size=18,
        color_hex=theme["secondary"],
    )

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, "#FFFDF9")
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Cooling Effect by District",
        0.8,
        0.7,
        5.2,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    editor.add_chart(
        slide_index,
        "bar_clustered",
        {
            "categories": ["North", "Central", "South", "West"],
            "series": [
                {"name": "Tree cover %", "values": [32, 24, 18, 28]},
                {"name": "Temp reduction", "values": [5.1, 3.8, 2.6, 4.4]},
            ],
        },
        0.9,
        1.6,
        7.4,
        4.0,
        style={
            "title": "Canopy coverage and cooling effect",
            "title_font_name": theme["font"],
            "title_font_size_pt": 18,
            "title_bold": True,
            "title_color_hex": theme["primary"],
            "legend_font_name": theme["font"],
            "legend_font_size_pt": 11,
            "legend_color_hex": theme["secondary"],
            "axis_font_name": theme["font"],
            "axis_font_size_pt": 11,
            "axis_color_hex": theme["secondary"],
            "series_colors": [theme["accent"], "#A3B18A"],
        },
    )
    add_text_block(
        editor,
        slide_index,
        "Higher canopy density consistently corresponds to lower peak pavement temperature.",
        8.6,
        2.0,
        3.0,
        2.0,
        font_name=theme["font"],
        font_size=16,
        color_hex=theme["secondary"],
    )

    slide_index = editor.add_slide(6)
    slide = editor.prs.slides[slide_index]
    apply_slide_background(slide, "#FFFDF9")
    add_accent_bar(slide, theme["accent"])
    add_text_block(
        editor,
        slide_index,
        "Field Sample Summary",
        0.8,
        0.7,
        4.5,
        0.6,
        font_name=theme["font"],
        font_size=22,
        color_hex=theme["primary"],
        bold=True,
    )
    editor.add_table(
        slide_index,
        4,
        3,
        [
            ["District", "Sites", "Avg canopy"],
            ["North", "18", "32%"],
            ["Central", "21", "24%"],
            ["South", "16", "18%"],
        ],
        0.9,
        1.6,
        7.6,
        2.5,
        style={
            "header_fill_hex": theme["accent"],
            "body_fill_hex": "#EFE7DA",
            "header_font_name": theme["font"],
            "body_font_name": theme["font"],
            "header_font_size_pt": 14,
            "body_font_size_pt": 12,
            "header_font_color_hex": "#FFFDF9",
            "body_font_color_hex": theme["primary"],
        },
    )
    editor.add_citation(slide_index, "Source: Urban Ecology Lab, field survey dataset")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    editor.prs.save(output_path)


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    output_paths = [
        OUTPUT_DIR / "presentation_modern_business.pptx",
        OUTPUT_DIR / "presentation_dark_analytics.pptx",
        OUTPUT_DIR / "presentation_academic_report.pptx",
    ]

    build_modern_business_deck(output_paths[0])
    build_dark_analytics_deck(output_paths[1])
    build_academic_report_deck(output_paths[2])

    for path in output_paths:
        print(path.relative_to(ROOT))


if __name__ == "__main__":
    main()
