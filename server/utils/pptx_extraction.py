from __future__ import annotations

import hashlib
import re
from collections import Counter
from dataclasses import dataclass
from io import BytesIO
from typing import Any

from pptx import Presentation as load_presentation
from pptx.presentation import Presentation as PptxPresentation

from server.utils.pptx_functions import PptxEditor
from server.utils.reward_models import (
    ExtractedChart,
    ExtractedImage,
    ExtractedPresentation,
    ExtractedShape,
    ExtractedSlide,
    ExtractedTable,
    ExtractedTextBlock,
)


_CITATION_PATTERN = re.compile(
    r"^\s*(source|sources|citation|citations|reference|references)\s*[:\-]",
    re.IGNORECASE,
)


@dataclass(slots=True)
class OpenPptxResult:
    presentation: PptxPresentation
    inspection_mode: str
    source_path: str | None = None


def open_presentation(
    presentation: PptxEditor | PptxPresentation | str,
) -> OpenPptxResult:
    if isinstance(presentation, PptxEditor):
        return OpenPptxResult(
            presentation=presentation.prs,
            inspection_mode="pptx_editor",
        )
    if isinstance(presentation, PptxPresentation):
        return OpenPptxResult(
            presentation=presentation,
            inspection_mode="presentation_object",
        )
    if isinstance(presentation, str):
        return OpenPptxResult(
            presentation=load_presentation(presentation),
            inspection_mode="pptx_path",
            source_path=presentation,
        )
    raise TypeError(
        "presentation must be a PptxEditor, pptx.Presentation, or .pptx path"
    )


def _safe_rgb(color_format: Any) -> str | None:
    try:
        rgb = color_format.rgb
    except Exception:
        return None
    return str(rgb) if rgb is not None else None


def _emu_to_inches(value: int | None) -> float:
    if value is None:
        return 0.0
    return float(value) / 914400.0


def _has_image(shape: Any) -> bool:
    try:
        shape.image
    except Exception:
        return False
    return True


def _shape_fill_hex(shape: Any) -> str | None:
    try:
        fill = shape.fill
    except Exception:
        return None
    try:
        return _safe_rgb(fill.fore_color)
    except Exception:
        return None


def _shape_line_hex(shape: Any) -> str | None:
    try:
        line = shape.line
    except Exception:
        return None
    try:
        return _safe_rgb(line.color)
    except Exception:
        return None


def _background_hex(slide: Any) -> str | None:
    try:
        return _safe_rgb(slide.background.fill.fore_color)
    except Exception:
        return None


def _infer_text_style(
    paragraph: Any,
) -> tuple[str | None, float | None, bool | None, bool | None, str | None]:
    candidates = []
    for run in getattr(paragraph, "runs", []):
        candidates.append(run.font)
    candidates.append(paragraph.font)

    font_name = None
    font_size = None
    bold = None
    italic = None
    color_hex = None

    for font in candidates:
        if font_name is None and getattr(font, "name", None) is not None:
            font_name = font.name
        if font_size is None and getattr(font, "size", None) is not None:
            font_size = float(font.size.pt)
        if bold is None and getattr(font, "bold", None) is not None:
            bold = bool(font.bold)
        if italic is None and getattr(font, "italic", None) is not None:
            italic = bool(font.italic)
        if color_hex is None:
            color_hex = _safe_rgb(getattr(font, "color", None))

    return font_name, font_size, bold, italic, color_hex


def _extract_text_block(shape: Any) -> ExtractedTextBlock:
    block = ExtractedTextBlock()
    for paragraph in shape.text_frame.paragraphs:
        font_name, font_size, bold, italic, color_hex = _infer_text_style(paragraph)
        block.paragraph_texts.append(paragraph.text)
        block.bullet_levels.append(getattr(paragraph, "level", None))
        block.font_names.append(font_name)
        block.font_sizes_pt.append(font_size)
        block.bold_flags.append(bold)
        block.italic_flags.append(italic)
        block.color_hexes.append(color_hex)
    return block


def _chart_type_name(chart: Any) -> str:
    chart_type = getattr(chart, "chart_type", None)
    if chart_type is None:
        return "unknown"
    return getattr(chart_type, "name", str(chart_type)).lower()


def _extract_chart(shape: Any) -> ExtractedChart:
    chart = shape.chart
    categories: list[str] = []
    series_payload: list[dict[str, Any]] = []
    axis_labels: dict[str, Any] = {}

    try:
        for category in chart.plots[0].categories:
            categories.append(str(getattr(category, "label", category)))
    except Exception:
        categories = []

    for series in chart.series:
        values: list[Any] = []
        try:
            values = list(series.values)
        except Exception:
            values = []
        series_payload.append(
            {
                "name": getattr(series, "name", None),
                "values": values,
            }
        )

    try:
        if chart.category_axis.has_title:
            axis_labels["category"] = chart.category_axis.axis_title.text_frame.text
    except Exception:
        pass
    try:
        if chart.value_axis.has_title:
            axis_labels["value"] = chart.value_axis.axis_title.text_frame.text
    except Exception:
        pass

    title = None
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text
    except Exception:
        title = None

    series_colors = []
    for series in chart.series:
        try:
            series_colors.append(_safe_rgb(series.format.fill.fore_color))
        except Exception:
            series_colors.append(None)

    return ExtractedChart(
        chart_type=_chart_type_name(chart),
        title=title,
        categories=categories,
        series=series_payload,
        has_legend=bool(getattr(chart, "has_legend", False)),
        axis_labels=axis_labels,
        style_metrics={"series_colors": series_colors},
    )


def _extract_table(shape: Any) -> ExtractedTable:
    table = shape.table
    cells: list[list[str]] = []
    for row in table.rows:
        cells.append([cell.text for cell in row.cells])
    return ExtractedTable(
        rows=len(table.rows),
        cols=len(table.columns),
        cells=cells,
        header_present=len(table.rows) > 1,
        style_metrics={},
    )


def _extract_image(shape: Any) -> ExtractedImage:
    image = shape.image
    width_px = None
    height_px = None
    try:
        width_px, height_px = image.size
    except Exception:
        pass

    content_hash = None
    try:
        content_hash = hashlib.sha256(image.blob).hexdigest()
    except Exception:
        content_hash = None

    return ExtractedImage(
        width_px=width_px,
        height_px=height_px,
        content_hash=content_hash,
        metadata={"ext": getattr(image, "ext", None)},
    )


def _shape_kind(shape: Any, slide_height_in: float) -> tuple[str, str | None]:
    if getattr(shape, "has_chart", False):
        return "chart", "chart"
    if getattr(shape, "has_table", False):
        return "table", "table"
    if _has_image(shape):
        return "image", "image"
    width = float(shape.width.inches)
    height = float(shape.height.inches)
    top = float(shape.top.inches)
    left = float(shape.left.inches)
    if top <= 0.3 and left <= 0.2 and width >= 10 and height <= 0.6:
        return "accent_bar", "accent_bar"
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text or ""
        if _CITATION_PATTERN.search(text) or (
            top >= max(slide_height_in - 1.2, 0) and height <= 0.7
        ):
            return "citation", "citation"
        return "text", None
    return "unknown", None


def _extract_shape(
    slide: Any,
    slide_index: int,
    shape: Any,
    z_index: int,
    *,
    slide_height_in: float,
) -> ExtractedShape:
    shape_kind, inferred_role = _shape_kind(shape, slide_height_in)
    raw_text = (
        shape.text_frame.text if getattr(shape, "has_text_frame", False) else None
    )
    text_blocks = []
    chart = None
    table = None
    image = None

    if getattr(shape, "has_text_frame", False):
        text_blocks.append(_extract_text_block(shape))
    if getattr(shape, "has_chart", False):
        chart = _extract_chart(shape)
    if getattr(shape, "has_table", False):
        table = _extract_table(shape)
    if _has_image(shape):
        image = _extract_image(shape)

    return ExtractedShape(
        shape_id=shape.shape_id,
        shape_kind=shape_kind,
        semantic_role=inferred_role,
        name=None,
        x=float(shape.left.inches),
        y=float(shape.top.inches),
        w=float(shape.width.inches),
        h=float(shape.height.inches),
        z_index=z_index,
        fill_color_hex=_shape_fill_hex(shape),
        line_color_hex=_shape_line_hex(shape),
        raw_text=raw_text,
        text_blocks=text_blocks,
        chart=chart,
        table=table,
        image=image,
        metadata={
            "native_name": getattr(shape, "name", None),
            "shape_type": str(getattr(shape, "shape_type", "unknown")),
            "slide_index": slide_index,
        },
    )


def _title_from_shapes(shapes: list[ExtractedShape]) -> str | None:
    text_shapes = [
        shape for shape in shapes if shape.shape_kind == "text" and shape.raw_text
    ]
    if not text_shapes:
        return None

    def score(shape: ExtractedShape) -> tuple[float, float, float]:
        font_sizes = [
            size
            for block in shape.text_blocks
            for size in block.font_sizes_pt
            if size is not None
        ]
        max_font = max(font_sizes, default=0.0)
        return (-shape.y, max_font, -len(shape.raw_text or ""))

    best_shape = sorted(text_shapes, key=score, reverse=True)[0]
    return best_shape.raw_text


def _shape_texts(shapes: list[ExtractedShape]) -> list[str]:
    texts = []
    for shape in shapes:
        if shape.raw_text:
            texts.append(shape.raw_text.strip())
        elif shape.table is not None:
            texts.extend(cell for row in shape.table.cells for cell in row if cell)
        elif shape.chart is not None:
            if shape.chart.title:
                texts.append(shape.chart.title)
            texts.extend(shape.chart.categories)
            for series in shape.chart.series:
                if series.get("name"):
                    texts.append(str(series["name"]))
                texts.extend(str(value) for value in series.get("values", []))
    return [text for text in texts if text]


def _font_metrics(shapes: list[ExtractedShape]) -> dict[str, Any]:
    font_sizes = [
        size
        for shape in shapes
        for block in shape.text_blocks
        for size in block.font_sizes_pt
        if size is not None
    ]
    font_names = [
        name
        for shape in shapes
        for block in shape.text_blocks
        for name in block.font_names
        if name
    ]
    return {
        "font_sizes_pt": font_sizes,
        "min_font_size_pt": min(font_sizes) if font_sizes else None,
        "max_font_size_pt": max(font_sizes) if font_sizes else None,
        "unique_font_families": sorted(set(font_names)),
        "unique_font_family_count": len(set(font_names)),
    }


def _layout_metrics(
    shapes: list[ExtractedShape],
    *,
    slide_width: float,
    slide_height: float,
) -> dict[str, Any]:
    slide_area = slide_width * slide_height if slide_width and slide_height else 0.0
    occupied_area = sum(shape.w * shape.h for shape in shapes)
    return {
        "shape_count": len(shapes),
        "occupied_area_ratio": (occupied_area / slide_area) if slide_area else 0.0,
        "text_shape_count": sum(
            1 for shape in shapes if shape.shape_kind in {"text", "citation"}
        ),
        "chart_count": sum(1 for shape in shapes if shape.shape_kind == "chart"),
        "table_count": sum(1 for shape in shapes if shape.shape_kind == "table"),
        "image_count": sum(1 for shape in shapes if shape.shape_kind == "image"),
    }


def _color_metrics(slide: Any, shapes: list[ExtractedShape]) -> dict[str, Any]:
    palette = [
        color
        for color in [
            _background_hex(slide),
            *[shape.fill_color_hex for shape in shapes],
            *[
                color_hex
                for shape in shapes
                for block in shape.text_blocks
                for color_hex in block.color_hexes
                if color_hex is not None
            ],
        ]
        if color
    ]
    return {
        "palette": sorted(set(palette)),
        "palette_counts": dict(Counter(palette)),
    }


def presentation_digest(presentation: PptxPresentation) -> str:
    payload = BytesIO()
    presentation.save(payload)
    return hashlib.sha256(payload.getvalue()).hexdigest()


class PptxExtractionService:
    def inspect_presentation(
        self,
        presentation: PptxEditor | PptxPresentation | str,
    ) -> ExtractedPresentation:
        opened = open_presentation(presentation)
        slides = []
        for slide_index, slide in enumerate(opened.presentation.slides, start=1):
            slides.append(
                self.inspect_slide(
                    slide_index,
                    presentation=opened.presentation,
                )
            )

        palette = Counter(
            color
            for slide in slides
            for color in slide.color_metrics.get("palette", [])
            if color
        )

        return ExtractedPresentation(
            slide_count=len(slides),
            slide_ids=[slide.slide_id for slide in slides],
            slides=slides,
            deck_metrics={
                "inspection_mode": opened.inspection_mode,
                "chart_count": sum(
                    slide.layout_metrics.get("chart_count", 0) for slide in slides
                ),
                "table_count": sum(
                    slide.layout_metrics.get("table_count", 0) for slide in slides
                ),
                "image_count": sum(
                    slide.layout_metrics.get("image_count", 0) for slide in slides
                ),
                "citation_count": sum(len(slide.citations) for slide in slides),
            },
            theme_summary={
                "dominant_colors": [color for color, _count in palette.most_common(5)],
            },
            metadata={
                "inspection_mode": opened.inspection_mode,
                "source_path": opened.source_path,
                "presentation_digest": presentation_digest(opened.presentation),
            },
        )

    def inspect_slide(
        self,
        slide_index: int,
        *,
        presentation: PptxEditor | PptxPresentation,
    ) -> ExtractedSlide:
        opened = open_presentation(presentation)
        if slide_index < 1 or slide_index > len(opened.presentation.slides):
            raise IndexError(f"Slide index {slide_index} out of range")

        slide = opened.presentation.slides[slide_index - 1]
        slide_width = _emu_to_inches(opened.presentation.slide_width)
        slide_height = _emu_to_inches(opened.presentation.slide_height)
        shapes = [
            _extract_shape(
                slide,
                slide_index,
                shape,
                z_index,
                slide_height_in=slide_height,
            )
            for z_index, shape in enumerate(slide.shapes)
        ]

        citations = [
            shape.raw_text
            for shape in shapes
            if shape.shape_kind == "citation" and shape.raw_text
        ]
        font_metrics = _font_metrics(shapes)
        title_text = None
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and getattr(
                title_shape, "has_text_frame", False
            ):
                title_text = title_shape.text_frame.text or None
        except Exception:
            title_text = None
        if not title_text:
            title_text = _title_from_shapes(shapes)

        return ExtractedSlide(
            slide_index=slide_index,
            slide_id=slide.slide_id,
            layout_name=getattr(slide.slide_layout, "name", None),
            background_color_hex=_background_hex(slide),
            title_text=title_text,
            all_text="\n".join(_shape_texts(shapes)).strip(),
            citations=[citation for citation in citations if citation],
            shapes=shapes,
            text_metrics={
                **font_metrics,
                "word_count": len(" ".join(_shape_texts(shapes)).split()),
            },
            layout_metrics=_layout_metrics(
                shapes,
                slide_width=slide_width,
                slide_height=slide_height,
            ),
            color_metrics=_color_metrics(slide, shapes),
            metadata={
                "inspection_mode": opened.inspection_mode,
            },
        )


__all__ = [
    "OpenPptxResult",
    "PptxExtractionService",
    "open_presentation",
    "presentation_digest",
]
