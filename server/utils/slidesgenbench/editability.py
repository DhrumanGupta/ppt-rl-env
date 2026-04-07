"""PEI (Presentation Editability Intelligence) taxonomy evaluation.

Implements the hierarchical knockout scoring from SlidesGen-Bench (Yang et al., 2026).
Levels L0-L5 are evaluated against the native PPTX object graph.
Failing a lower level precludes credit at higher ones (knockout rule).
"""

from __future__ import annotations

from typing import Any

from pptx.presentation import Presentation as PptxPresentation

from server.utils.pptx_functions import PptxEditor
from server.utils.pptx_extraction import open_presentation

PEI_LEVEL_REWARD: dict[int, float] = {
    0: 0.00,
    1: 0.20,
    2: 0.45,
    3: 0.70,
    4: 0.90,
    5: 1.00,
}


def _has_selectable_text(prs: PptxPresentation) -> bool:
    """L1 gate: at least one slide has a native text frame with selectable text."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text or ""
                if text.strip():
                    return True
    return False


def _text_is_not_fragmented(prs: PptxPresentation) -> bool:
    """L1 check: paragraphs are not split into single-line boxes.

    Heuristic: at least one text frame contains multiple paragraphs,
    indicating reflow capability rather than OCR-style reconstruction.
    """
    multi_para_count = 0
    single_line_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            paragraphs = shape.text_frame.paragraphs
            non_empty = [p for p in paragraphs if (p.text or "").strip()]
            if len(non_empty) > 1:
                multi_para_count += 1
            elif len(non_empty) == 1:
                single_line_count += 1
    if multi_para_count + single_line_count == 0:
        return False
    return multi_para_count > 0


def _has_vector_shapes(prs: PptxPresentation) -> bool:
    """L2 gate: presentation contains native vector shapes (not just images/text)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for slide in prs.slides:
        for shape in slide.shapes:
            shape_type = getattr(shape, "shape_type", None)
            if shape_type in (
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.FREEFORM,
            ):
                return True
            if getattr(shape, "has_chart", False):
                return True
            if getattr(shape, "has_table", False):
                return True
    return False


def _has_slide_master_inheritance(prs: PptxPresentation) -> bool:
    """L3 gate: slides inherit from slide master / slide layouts."""
    if not prs.slide_masters:
        return False
    if not prs.slide_layouts:
        return False
    for slide in prs.slides:
        layout = getattr(slide, "slide_layout", None)
        if layout is not None:
            return True
    return False


def _has_logical_grouping(prs: PptxPresentation) -> bool:
    """L3 check: at least one slide has grouped shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                return True
    return False


def _has_native_chart_data(prs: PptxPresentation) -> bool:
    """L4 gate: at least one chart is a native <c:chart> with data, not a static image."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            chart = shape.chart
            try:
                if chart.series and len(list(chart.series)) > 0:
                    first_series = list(chart.series)[0]
                    values = list(first_series.values)
                    if values:
                        return True
            except Exception:
                continue
    return False


def _has_animation_or_media(prs: PptxPresentation) -> bool:
    """L5 gate: slides contain animation sequences or embedded media."""
    anim_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    for slide in prs.slides:
        slide_xml = slide._element
        timing = slide_xml.findall(f"{{{anim_ns}}}timing")
        if timing:
            for t in timing:
                seq = t.findall(".//" + f"{{{anim_ns}}}seq")
                if seq:
                    return True
        transition = slide_xml.findall(f"{{{anim_ns}}}transition")
        if transition:
            return True
        for shape in slide.shapes:
            if getattr(shape, "media_type", None) is not None:
                return True
    return False


def evaluate_pei_level(
    presentation: PptxEditor | PptxPresentation | str,
) -> dict[str, Any]:
    """Evaluate the PEI level of a presentation using the knockout rule.

    Returns a dict with:
        pei_level: int (0-5)
        pei_reward: float (0.0-1.0)
        gate_results: dict mapping each gate to pass/fail
    """
    opened = open_presentation(presentation)
    prs = opened.presentation

    if not prs.slides or len(prs.slides) == 0:
        return {
            "pei_level": 0,
            "pei_reward": PEI_LEVEL_REWARD[0],
            "gate_results": {"L1_text": False},
        }

    gate_results: dict[str, bool] = {}

    # L1: Text Integrity
    has_text = _has_selectable_text(prs)
    not_fragmented = _text_is_not_fragmented(prs) if has_text else False
    gate_results["L1_selectable_text"] = has_text
    gate_results["L1_not_fragmented"] = not_fragmented
    if not has_text or not not_fragmented:
        level = 1 if has_text else 0
        return {
            "pei_level": level,
            "pei_reward": PEI_LEVEL_REWARD[level],
            "gate_results": gate_results,
        }

    # L2: Vector Fidelity
    has_vectors = _has_vector_shapes(prs)
    gate_results["L2_vector_shapes"] = has_vectors
    if not has_vectors:
        return {
            "pei_level": 1,
            "pei_reward": PEI_LEVEL_REWARD[1],
            "gate_results": gate_results,
        }

    # L3: Structural Logic
    has_master = _has_slide_master_inheritance(prs)
    has_groups = _has_logical_grouping(prs)
    gate_results["L3_master_inheritance"] = has_master
    gate_results["L3_logical_grouping"] = has_groups
    if not has_master:
        return {
            "pei_level": 2,
            "pei_reward": PEI_LEVEL_REWARD[2],
            "gate_results": gate_results,
        }

    # L4: Parametric (Native Data)
    has_native_charts = _has_native_chart_data(prs)
    gate_results["L4_native_chart_data"] = has_native_charts
    if not has_native_charts:
        return {
            "pei_level": 3,
            "pei_reward": PEI_LEVEL_REWARD[3],
            "gate_results": gate_results,
        }

    # L5: Cinematic (Animations + Media)
    try:
        has_anim = _has_animation_or_media(prs)
    except Exception:
        has_anim = False
    gate_results["L5_animation_or_media"] = has_anim
    if not has_anim:
        return {
            "pei_level": 4,
            "pei_reward": PEI_LEVEL_REWARD[4],
            "gate_results": gate_results,
        }

    return {
        "pei_level": 5,
        "pei_reward": PEI_LEVEL_REWARD[5],
        "gate_results": gate_results,
    }


__all__ = ["PEI_LEVEL_REWARD", "evaluate_pei_level"]
