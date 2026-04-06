import os
from copy import deepcopy
from dataclasses import dataclass
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


@dataclass
class ThemeBinding:
    apply_method: str
    slide_id: int
    shape_id: Optional[int]
    raw_kwargs: Dict[str, Any]


# a bit of own, a bit of vibe, makes everyone happy
class PptxEditor:
    def __init__(self, path: Optional[str] = None):
        self.prs = Presentation(path) if path else Presentation()
        self._themes: Dict[str, Dict[str, Any]] = {}
        self._theme_bindings: Dict[str, List[ThemeBinding]] = {}

    def _validate_slide(self, index: int):
        if index < 0 or index >= len(self.prs.slides):
            raise IndexError(f"Slide index {index} out of range")

    def _get_slide_by_id(self, slide_id: int):
        for slide in self.prs.slides:
            if slide.slide_id == slide_id:
                return slide
        raise KeyError(f"Slide id {slide_id} not found")

    def get_slide_id(self, slide_index: int) -> int:
        self._validate_slide(slide_index)
        return self.prs.slides[slide_index].slide_id

    def get_slide_index(self, slide_id: int) -> int:
        for index, slide in enumerate(self.prs.slides):
            if slide.slide_id == slide_id:
                return index
        raise KeyError(f"Slide id {slide_id} not found")

    def _get_shape(self, slide_index: int, shape_index: int):
        self._validate_slide(slide_index)
        shapes = list(self.prs.slides[slide_index].shapes)
        if shape_index < 0 or shape_index >= len(shapes):
            raise IndexError(f"Shape index {shape_index} out of range")
        return shapes[shape_index]

    def _get_shape_by_id(self, slide_id: int, shape_id: int):
        slide = self._get_slide_by_id(slide_id)
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
        raise KeyError(f"Shape id {shape_id} not found on slide {slide_id}")

    def get_shape_id(self, slide_index: int, shape_index: int) -> int:
        return self._get_shape(slide_index, shape_index).shape_id

    def get_shape_index(self, slide_id: int, shape_id: int) -> int:
        slide = self._get_slide_by_id(slide_id)
        for index, shape in enumerate(slide.shapes):
            if shape.shape_id == shape_id:
                return index
        raise KeyError(f"Shape id {shape_id} not found on slide {slide_id}")

    def _rgb_color(self, color_hex: str) -> RGBColor:
        return RGBColor.from_string(color_hex.replace("#", ""))

    def _is_theme_token(self, value: Any) -> bool:
        return isinstance(value, str) and value.startswith("<") and value.endswith(">")

    def _extract_theme_token(self, value: str) -> str:
        return value[1:-1]

    def _contains_theme_tokens(self, payload: Any) -> bool:
        if self._is_theme_token(payload):
            return True
        if isinstance(payload, dict):
            return any(self._contains_theme_tokens(value) for value in payload.values())
        if isinstance(payload, (list, tuple)):
            return any(self._contains_theme_tokens(value) for value in payload)
        return False

    def _get_registered_theme(self, name: str) -> Dict[str, Any]:
        if name not in self._themes:
            raise ValueError(f"Theme '{name}' is not registered")
        return self._themes[name]

    def _resolve_theme_token(
        self, theme_name: str, token_name: str, seen: Optional[set[str]] = None
    ) -> Any:
        theme = self._get_registered_theme(theme_name)
        if token_name not in theme:
            raise ValueError(f"Theme '{theme_name}' has no token '{token_name}'")

        seen = seen or set()
        if token_name in seen:
            raise ValueError(f"Cyclic theme token reference for '{token_name}'")

        return self._resolve_theme_payload(
            theme_name, theme[token_name], seen | {token_name}
        )

    def _resolve_theme_payload(
        self, theme_name: str, payload: Any, seen: Optional[set[str]] = None
    ) -> Any:
        if self._is_theme_token(payload):
            return self._resolve_theme_token(
                theme_name, self._extract_theme_token(payload), seen
            )
        if isinstance(payload, dict):
            return {
                key: self._resolve_theme_payload(theme_name, value, seen)
                for key, value in payload.items()
            }
        if isinstance(payload, list):
            return [
                self._resolve_theme_payload(theme_name, value, seen)
                for value in payload
            ]
        if isinstance(payload, tuple):
            return tuple(
                self._resolve_theme_payload(theme_name, value, seen)
                for value in payload
            )
        return payload

    def _resolve_bound_kwargs(
        self, theme_name: Optional[str], raw_kwargs: Dict[str, Any]
    ) -> Dict[str, Any]:
        if not self._contains_theme_tokens(raw_kwargs):
            return deepcopy(raw_kwargs)
        if theme_name is None:
            raise ValueError("theme_name is required when using theme tokens")

        resolved_kwargs = self._resolve_theme_payload(theme_name, deepcopy(raw_kwargs))
        if not isinstance(resolved_kwargs, dict):
            raise ValueError("Resolved theme payload must be a dictionary")
        return resolved_kwargs

    def _bind_theme(
        self,
        theme_name: Optional[str],
        apply_method: str,
        slide_id: int,
        shape_id: Optional[int],
        raw_kwargs: Dict[str, Any],
        bind_theme: bool,
    ):
        if not bind_theme or not self._contains_theme_tokens(raw_kwargs):
            return
        if theme_name is None:
            raise ValueError("theme_name is required when using theme tokens")

        self._get_registered_theme(theme_name)
        self._theme_bindings.setdefault(theme_name, []).append(
            ThemeBinding(
                apply_method=apply_method,
                slide_id=slide_id,
                shape_id=shape_id,
                raw_kwargs=deepcopy(raw_kwargs),
            )
        )

    def _reapply_binding(self, theme_name: str, binding: ThemeBinding) -> bool:
        method = getattr(self, binding.apply_method, None)
        if method is None:
            return False

        resolved_kwargs = self._resolve_bound_kwargs(theme_name, binding.raw_kwargs)

        try:
            if binding.shape_id is None:
                method(binding.slide_id, bind_theme=False, **resolved_kwargs)
            else:
                method(
                    binding.slide_id,
                    binding.shape_id,
                    bind_theme=False,
                    **resolved_kwargs,
                )
        except KeyError:
            return False

        return True

    def _reapply_theme(self, theme_name: str):
        bindings = self._theme_bindings.get(theme_name, [])
        active_bindings = []

        for binding in bindings:
            if self._reapply_binding(theme_name, binding):
                active_bindings.append(binding)

        self._theme_bindings[theme_name] = active_bindings

    def register_theme(self, name: str, theme: Dict[str, Any]) -> str:
        if not name:
            raise ValueError("Theme name cannot be empty")

        had_existing_theme = name in self._themes
        self._themes[name] = deepcopy(theme)
        self._theme_bindings.setdefault(name, [])

        if had_existing_theme and self._theme_bindings[name]:
            self._reapply_theme(name)

        return name

    def get_theme(self, name: str = "default") -> Dict[str, Any]:
        return deepcopy(self._get_registered_theme(name))

    def update_theme(self, name: str, updates: Dict[str, Any]):
        theme = self._get_registered_theme(name)
        theme.update(deepcopy(updates))
        self._reapply_theme(name)

    def _apply_font_style(
        self,
        font,
        font_name: Optional[str] = None,
        font_size_pt: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color_hex: Optional[str] = None,
    ):
        if font_name is not None:
            font.name = font_name
        if font_size_pt is not None:
            font.size = Pt(font_size_pt)
        if bold is not None:
            font.bold = bold
        if italic is not None:
            font.italic = italic
        if color_hex is not None:
            font.color.rgb = self._rgb_color(color_hex)

    def _apply_text_frame_style(
        self,
        text_frame,
        font_name: Optional[str] = None,
        font_size_pt: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color_hex: Optional[str] = None,
    ):
        for paragraph in text_frame.paragraphs:
            self._apply_font_style(
                paragraph.font,
                font_name=font_name,
                font_size_pt=font_size_pt,
                bold=bold,
                italic=italic,
                color_hex=color_hex,
            )
            for run in paragraph.runs:
                self._apply_font_style(
                    run.font,
                    font_name=font_name,
                    font_size_pt=font_size_pt,
                    bold=bold,
                    italic=italic,
                    color_hex=color_hex,
                )

    def add_slide(self, layout_index: int = 1) -> int:
        layout_index = max(0, min(layout_index, len(self.prs.slide_layouts) - 1))
        layout = self.prs.slide_layouts[layout_index]
        self.prs.slides.add_slide(layout)
        return len(self.prs.slides) - 1

    def reorder_slide(self, old_index: int, new_index: int):
        self._validate_slide(old_index)
        xml_slides = self.prs.slides._sldIdLst
        if new_index < 0 or new_index >= len(xml_slides):
            raise IndexError("new_index out of bounds")
        slide_elem = xml_slides[old_index]
        del xml_slides[old_index]
        xml_slides.insert(new_index, slide_elem)

    def set_slide_layout(self, slide_index: int, layout_index: int):
        self._validate_slide(slide_index)
        if layout_index < 0 or layout_index >= len(self.prs.slide_layouts):
            raise IndexError("Invalid layout index")
        slide = self.prs.slides[slide_index]
        new_layout = self.prs.slide_layouts[layout_index]

        layout_rel = None
        for rel in slide.part.rels.values():
            if "slideLayout" in rel.reltype:
                layout_rel = rel
                break

        if layout_rel:
            slide.part.drop_rel(layout_rel.rId)

        slide.part.relate_to(
            new_layout.part,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        )
        if hasattr(slide, "_slide_layout"):
            del slide._slide_layout

    def set_slide_background_by_id(
        self,
        slide_id: int,
        color_hex: str,
        *,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ):
        raw_kwargs = {"color_hex": color_hex}
        resolved_kwargs = self._resolve_bound_kwargs(theme_name, raw_kwargs)
        slide = self._get_slide_by_id(slide_id)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb_color(resolved_kwargs["color_hex"])
        self._bind_theme(
            theme_name,
            "set_slide_background_by_id",
            slide_id,
            None,
            raw_kwargs,
            bind_theme,
        )

    def style_shape_fill_by_id(
        self,
        slide_id: int,
        shape_id: int,
        fill_color_hex: Optional[str] = None,
        *,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ):
        raw_kwargs = {"fill_color_hex": fill_color_hex}
        resolved_kwargs = self._resolve_bound_kwargs(theme_name, raw_kwargs)
        shape = self._get_shape_by_id(slide_id, shape_id)
        if resolved_kwargs["fill_color_hex"] is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb_color(
                resolved_kwargs["fill_color_hex"]
            )
        self._bind_theme(
            theme_name,
            "style_shape_fill_by_id",
            slide_id,
            shape_id,
            raw_kwargs,
            bind_theme,
        )

    def add_accent_bar_by_id(
        self,
        slide_id: int,
        color_hex: str,
        *,
        height: float = 0.35,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ) -> int:
        slide = self._get_slide_by_id(slide_id)
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.33),
            Inches(height),
        )
        accent.line.fill.background()
        shape_id = accent.shape_id
        self.style_shape_fill_by_id(
            slide_id,
            shape_id,
            fill_color_hex=color_hex,
            theme_name=theme_name,
            bind_theme=bind_theme,
        )
        return shape_id

    def add_textbox_by_id(
        self, slide_id: int, x: float, y: float, cx: float, cy: float
    ) -> int:
        slide = self._get_slide_by_id(slide_id)
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cx), Inches(cy))
        return textbox.shape_id

    def insert_text_by_id(self, slide_id: int, shape_id: int, text: str):
        shape = self._get_shape_by_id(slide_id, shape_id)
        if not shape.has_text_frame:
            raise ValueError("Shape has no text frame.")
        shape.text_frame.text = str(text)

    def insert_text(self, slide_index: int, shape_index: int, text: str):
        self.insert_text_by_id(
            self.get_slide_id(slide_index),
            self.get_shape_id(slide_index, shape_index),
            text,
        )

    def add_chart_by_id(
        self,
        slide_id: int,
        chart_type_str: str,
        chart_data_dict: Dict[str, Any],
        x: float,
        y: float,
        cx: float,
        cy: float,
        style: Optional[Dict[str, Any]] = None,
        *,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ) -> int:
        slide = self._get_slide_by_id(slide_id)

        chart_data = CategoryChartData()
        chart_data.categories = chart_data_dict.get("categories", [])
        for series in chart_data_dict.get("series", []):
            chart_data.add_series(series["name"], tuple(series["values"]))

        chart_type_enum = getattr(
            XL_CHART_TYPE, chart_type_str.upper(), XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        chart = slide.shapes.add_chart(
            chart_type_enum, Inches(x), Inches(y), Inches(cx), Inches(cy), chart_data
        )
        shape_id = chart.shape_id
        if style:
            self.style_chart_by_id(
                slide_id,
                shape_id,
                theme_name=theme_name,
                bind_theme=bind_theme,
                **style,
            )
        return shape_id

    def add_chart(
        self,
        slide_index: int,
        chart_type_str: str,
        chart_data_dict: Dict[str, Any],
        x: float,
        y: float,
        cx: float,
        cy: float,
        style: Optional[Dict[str, Any]] = None,
    ):
        return self.add_chart_by_id(
            self.get_slide_id(slide_index),
            chart_type_str,
            chart_data_dict,
            x,
            y,
            cx,
            cy,
            style=style,
        )

    def add_table_by_id(
        self,
        slide_id: int,
        rows: int,
        cols: int,
        table_data: List[List[str]],
        x: float,
        y: float,
        cx: float,
        cy: float,
        style: Optional[Dict[str, Any]] = None,
        *,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ) -> int:
        slide = self._get_slide_by_id(slide_id)
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(x), Inches(y), Inches(cx), Inches(cy)
        )
        table = table_shape.table

        for r in range(min(rows, len(table_data))):
            for c in range(min(cols, len(table_data[r]))):
                table.cell(r, c).text = str(table_data[r][c])

        shape_id = table_shape.shape_id
        if style:
            self.style_table_by_id(
                slide_id,
                shape_id,
                theme_name=theme_name,
                bind_theme=bind_theme,
                **style,
            )
        return shape_id

    def add_table(
        self,
        slide_index: int,
        rows: int,
        cols: int,
        table_data: List[List[str]],
        x: float,
        y: float,
        cx: float,
        cy: float,
        style: Optional[Dict[str, Any]] = None,
    ):
        return self.add_table_by_id(
            self.get_slide_id(slide_index),
            rows,
            cols,
            table_data,
            x,
            y,
            cx,
            cy,
            style=style,
        )

    def add_image_by_id(
        self,
        slide_id: int,
        image_path: str,
        x: float,
        y: float,
        cx: Optional[float] = None,
        cy: Optional[float] = None,
    ) -> int:
        if not os.path.isfile(image_path):
            raise FileNotFoundError("Image path invalid.")

        slide = self._get_slide_by_id(slide_id)
        kwargs = {"left": Inches(x), "top": Inches(y)}
        if cx is not None:
            kwargs["width"] = Inches(cx)
        if cy is not None:
            kwargs["height"] = Inches(cy)

        picture = slide.shapes.add_picture(image_path, **kwargs)
        return picture.shape_id

    def add_image(
        self,
        slide_index: int,
        image_path: str,
        x: float,
        y: float,
        cx: Optional[float] = None,
        cy: Optional[float] = None,
    ):
        return self.add_image_by_id(
            self.get_slide_id(slide_index), image_path, x, y, cx=cx, cy=cy
        )

    def add_citation_by_id(
        self,
        slide_id: int,
        text: str,
        style: Optional[Dict[str, Any]] = None,
        *,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ) -> int:
        shape_id = self.add_textbox_by_id(slide_id, 0.5, 6.8, 9.0, 0.4)
        self.insert_text_by_id(slide_id, shape_id, text)

        citation_style = {"font_size_pt": 10, "italic": True}
        if style:
            citation_style.update(style)

        self.style_text_by_id(
            slide_id,
            shape_id,
            theme_name=theme_name,
            bind_theme=bind_theme,
            **citation_style,
        )
        return shape_id

    def add_citation(self, slide_index: int, text: str):
        return self.add_citation_by_id(self.get_slide_id(slide_index), text)

    def style_text_by_id(
        self,
        slide_id: int,
        shape_id: int,
        font_name: Optional[Any] = None,
        font_size_pt: Optional[Any] = None,
        bold: Optional[Any] = None,
        italic: Optional[Any] = None,
        color_hex: Optional[Any] = None,
        *,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ):
        raw_kwargs = {
            "font_name": font_name,
            "font_size_pt": font_size_pt,
            "bold": bold,
            "italic": italic,
            "color_hex": color_hex,
        }
        resolved_kwargs = self._resolve_bound_kwargs(theme_name, raw_kwargs)

        shape = self._get_shape_by_id(slide_id, shape_id)
        if not shape.has_text_frame:
            raise ValueError("Shape has no text frame.")

        self._apply_text_frame_style(shape.text_frame, **resolved_kwargs)
        self._bind_theme(
            theme_name,
            "style_text_by_id",
            slide_id,
            shape_id,
            raw_kwargs,
            bind_theme,
        )

    def style_update(
        self,
        slide_index: int,
        shape_index: int,
        font_name: Optional[str] = None,
        font_size_pt: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color_hex: Optional[str] = None,
    ):
        self.style_text_by_id(
            self.get_slide_id(slide_index),
            self.get_shape_id(slide_index, shape_index),
            font_name=font_name,
            font_size_pt=font_size_pt,
            bold=bold,
            italic=italic,
            color_hex=color_hex,
        )

    def style_table_by_id(
        self,
        slide_id: int,
        shape_id: int,
        header_fill_hex: Optional[Any] = None,
        body_fill_hex: Optional[Any] = None,
        header_font_name: Optional[Any] = None,
        body_font_name: Optional[Any] = None,
        header_font_size_pt: Optional[Any] = None,
        body_font_size_pt: Optional[Any] = None,
        header_bold: Optional[Any] = True,
        body_bold: Optional[Any] = None,
        header_italic: Optional[Any] = None,
        body_italic: Optional[Any] = None,
        header_font_color_hex: Optional[Any] = None,
        body_font_color_hex: Optional[Any] = None,
        *,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ):
        raw_kwargs = {
            "header_fill_hex": header_fill_hex,
            "body_fill_hex": body_fill_hex,
            "header_font_name": header_font_name,
            "body_font_name": body_font_name,
            "header_font_size_pt": header_font_size_pt,
            "body_font_size_pt": body_font_size_pt,
            "header_bold": header_bold,
            "body_bold": body_bold,
            "header_italic": header_italic,
            "body_italic": body_italic,
            "header_font_color_hex": header_font_color_hex,
            "body_font_color_hex": body_font_color_hex,
        }
        resolved_kwargs = self._resolve_bound_kwargs(theme_name, raw_kwargs)

        shape = self._get_shape_by_id(slide_id, shape_id)
        if not shape.has_table:
            raise ValueError("Shape has no table.")

        table = shape.table
        for row_idx, row in enumerate(table.rows):
            is_header = row_idx == 0
            fill_hex = (
                resolved_kwargs["header_fill_hex"]
                if is_header
                else resolved_kwargs["body_fill_hex"]
            )
            font_name = (
                resolved_kwargs["header_font_name"]
                if is_header
                else resolved_kwargs["body_font_name"]
            )
            font_size_pt = (
                resolved_kwargs["header_font_size_pt"]
                if is_header
                else resolved_kwargs["body_font_size_pt"]
            )
            bold = (
                resolved_kwargs["header_bold"]
                if is_header
                else resolved_kwargs["body_bold"]
            )
            italic = (
                resolved_kwargs["header_italic"]
                if is_header
                else resolved_kwargs["body_italic"]
            )
            font_color_hex = (
                resolved_kwargs["header_font_color_hex"]
                if is_header
                else resolved_kwargs["body_font_color_hex"]
            )

            for cell in row.cells:
                if fill_hex is not None:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._rgb_color(fill_hex)
                self._apply_text_frame_style(
                    cell.text_frame,
                    font_name=font_name,
                    font_size_pt=font_size_pt,
                    bold=bold,
                    italic=italic,
                    color_hex=font_color_hex,
                )

        self._bind_theme(
            theme_name,
            "style_table_by_id",
            slide_id,
            shape_id,
            raw_kwargs,
            bind_theme,
        )

    def style_table(
        self,
        slide_index: int,
        shape_index: int,
        header_fill_hex: Optional[str] = None,
        body_fill_hex: Optional[str] = None,
        header_font_name: Optional[str] = None,
        body_font_name: Optional[str] = None,
        header_font_size_pt: Optional[int] = None,
        body_font_size_pt: Optional[int] = None,
        header_bold: Optional[bool] = True,
        body_bold: Optional[bool] = None,
        header_italic: Optional[bool] = None,
        body_italic: Optional[bool] = None,
        header_font_color_hex: Optional[str] = None,
        body_font_color_hex: Optional[str] = None,
    ):
        self.style_table_by_id(
            self.get_slide_id(slide_index),
            self.get_shape_id(slide_index, shape_index),
            header_fill_hex=header_fill_hex,
            body_fill_hex=body_fill_hex,
            header_font_name=header_font_name,
            body_font_name=body_font_name,
            header_font_size_pt=header_font_size_pt,
            body_font_size_pt=body_font_size_pt,
            header_bold=header_bold,
            body_bold=body_bold,
            header_italic=header_italic,
            body_italic=body_italic,
            header_font_color_hex=header_font_color_hex,
            body_font_color_hex=body_font_color_hex,
        )

    def style_chart_by_id(
        self,
        slide_id: int,
        shape_id: int,
        title: Optional[Any] = None,
        title_font_name: Optional[Any] = None,
        title_font_size_pt: Optional[Any] = None,
        title_bold: Optional[Any] = None,
        title_italic: Optional[Any] = None,
        title_color_hex: Optional[Any] = None,
        legend_font_name: Optional[Any] = None,
        legend_font_size_pt: Optional[Any] = None,
        legend_bold: Optional[Any] = None,
        legend_italic: Optional[Any] = None,
        legend_color_hex: Optional[Any] = None,
        axis_font_name: Optional[Any] = None,
        axis_font_size_pt: Optional[Any] = None,
        axis_bold: Optional[Any] = None,
        axis_italic: Optional[Any] = None,
        axis_color_hex: Optional[Any] = None,
        series_colors: Optional[List[Any]] = None,
        *,
        theme_name: Optional[str] = None,
        bind_theme: bool = False,
    ):
        raw_kwargs = {
            "title": title,
            "title_font_name": title_font_name,
            "title_font_size_pt": title_font_size_pt,
            "title_bold": title_bold,
            "title_italic": title_italic,
            "title_color_hex": title_color_hex,
            "legend_font_name": legend_font_name,
            "legend_font_size_pt": legend_font_size_pt,
            "legend_bold": legend_bold,
            "legend_italic": legend_italic,
            "legend_color_hex": legend_color_hex,
            "axis_font_name": axis_font_name,
            "axis_font_size_pt": axis_font_size_pt,
            "axis_bold": axis_bold,
            "axis_italic": axis_italic,
            "axis_color_hex": axis_color_hex,
            "series_colors": series_colors,
        }
        resolved_kwargs = self._resolve_bound_kwargs(theme_name, raw_kwargs)

        shape = self._get_shape_by_id(slide_id, shape_id)
        if not shape.has_chart:
            raise ValueError("Shape has no chart.")

        chart = shape.chart

        if resolved_kwargs["title"] is not None:
            chart.has_title = True
            chart.chart_title.text_frame.text = str(resolved_kwargs["title"])
        if chart.has_title:
            self._apply_text_frame_style(
                chart.chart_title.text_frame,
                font_name=resolved_kwargs["title_font_name"],
                font_size_pt=resolved_kwargs["title_font_size_pt"],
                bold=resolved_kwargs["title_bold"],
                italic=resolved_kwargs["title_italic"],
                color_hex=resolved_kwargs["title_color_hex"],
            )

        if chart.has_legend:
            self._apply_font_style(
                chart.legend.font,
                font_name=resolved_kwargs["legend_font_name"],
                font_size_pt=resolved_kwargs["legend_font_size_pt"],
                bold=resolved_kwargs["legend_bold"],
                italic=resolved_kwargs["legend_italic"],
                color_hex=resolved_kwargs["legend_color_hex"],
            )

        for axis_name in ("category_axis", "value_axis"):
            axis = getattr(chart, axis_name, None)
            if axis is not None and hasattr(axis, "tick_labels"):
                self._apply_font_style(
                    axis.tick_labels.font,
                    font_name=resolved_kwargs["axis_font_name"],
                    font_size_pt=resolved_kwargs["axis_font_size_pt"],
                    bold=resolved_kwargs["axis_bold"],
                    italic=resolved_kwargs["axis_italic"],
                    color_hex=resolved_kwargs["axis_color_hex"],
                )

        if resolved_kwargs["series_colors"]:
            for series, color_hex in zip(
                chart.series, resolved_kwargs["series_colors"]
            ):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = self._rgb_color(color_hex)
                series.format.line.color.rgb = self._rgb_color(color_hex)

        self._bind_theme(
            theme_name,
            "style_chart_by_id",
            slide_id,
            shape_id,
            raw_kwargs,
            bind_theme,
        )

    def style_chart(
        self,
        slide_index: int,
        shape_index: int,
        title: Optional[str] = None,
        title_font_name: Optional[str] = None,
        title_font_size_pt: Optional[int] = None,
        title_bold: Optional[bool] = None,
        title_italic: Optional[bool] = None,
        title_color_hex: Optional[str] = None,
        legend_font_name: Optional[str] = None,
        legend_font_size_pt: Optional[int] = None,
        legend_bold: Optional[bool] = None,
        legend_italic: Optional[bool] = None,
        legend_color_hex: Optional[str] = None,
        axis_font_name: Optional[str] = None,
        axis_font_size_pt: Optional[int] = None,
        axis_bold: Optional[bool] = None,
        axis_italic: Optional[bool] = None,
        axis_color_hex: Optional[str] = None,
        series_colors: Optional[List[str]] = None,
    ):
        self.style_chart_by_id(
            self.get_slide_id(slide_index),
            self.get_shape_id(slide_index, shape_index),
            title=title,
            title_font_name=title_font_name,
            title_font_size_pt=title_font_size_pt,
            title_bold=title_bold,
            title_italic=title_italic,
            title_color_hex=title_color_hex,
            legend_font_name=legend_font_name,
            legend_font_size_pt=legend_font_size_pt,
            legend_bold=legend_bold,
            legend_italic=legend_italic,
            legend_color_hex=legend_color_hex,
            axis_font_name=axis_font_name,
            axis_font_size_pt=axis_font_size_pt,
            axis_bold=axis_bold,
            axis_italic=axis_italic,
            axis_color_hex=axis_color_hex,
            series_colors=series_colors,
        )


if __name__ == "__main__":
    pptx_editor = PptxEditor()
    slide_index = pptx_editor.add_slide(6)
    slide_id = pptx_editor.get_slide_id(slide_index)
    textbox_id = pptx_editor.add_textbox_by_id(slide_id, 0.5, 0.5, 9.0, 1.0)
    body_id = pptx_editor.add_textbox_by_id(slide_id, 0.5, 1.7, 9.0, 1.5)
    pptx_editor.insert_text_by_id(slide_id, textbox_id, "Hello, World!")
    pptx_editor.insert_text_by_id(slide_id, body_id, "Item 1\nItem 2\nItem 3")

    ppt = pptx_editor.prs
    os.makedirs("outputs", exist_ok=True)
    ppt.save("outputs/output.pptx")
