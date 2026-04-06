from typing import Any, Dict, Optional

from pptx.util import Inches

from server.utils.pptx_functions import PptxEditor


_SHAPE_ID_UNSET = object()


def register_theme(
    editor: PptxEditor, theme: Dict[str, Any], name: str = "default"
) -> str:
    return editor.register_theme(name, theme)


def update_theme(
    editor: PptxEditor, updates: Dict[str, Any], name: str = "default"
) -> None:
    editor.update_theme(name, updates)


def _remove_bindings(
    editor: PptxEditor,
    *,
    slide_id: Optional[int] = None,
    shape_id: Any = _SHAPE_ID_UNSET,
    apply_method: Optional[str] = None,
    theme_name: Optional[str] = None,
) -> None:
    theme_names = (
        [theme_name] if theme_name is not None else list(editor._theme_bindings)
    )

    for current_theme_name in theme_names:
        remaining = []
        for binding in editor._theme_bindings.get(current_theme_name, []):
            if slide_id is not None and binding.slide_id != slide_id:
                remaining.append(binding)
                continue
            if shape_id is not _SHAPE_ID_UNSET and binding.shape_id != shape_id:
                remaining.append(binding)
                continue
            if apply_method is not None and binding.apply_method != apply_method:
                remaining.append(binding)
                continue
        editor._theme_bindings[current_theme_name] = remaining


def _get_required(spec: Dict[str, Any], key: str) -> Any:
    if key not in spec:
        raise ValueError(f"Shape spec missing required field '{key}'")
    return spec[key]


def _get_optional_geometry(spec: Dict[str, Any]) -> Dict[str, float]:
    geometry = {}
    for source_key, target_key in (
        ("x", "left"),
        ("y", "top"),
        ("w", "width"),
        ("h", "height"),
    ):
        if source_key in spec:
            geometry[target_key] = spec[source_key]
    return geometry


def _apply_geometry_updates(
    editor: PptxEditor, slide_id: int, shape_id: int, spec: Dict[str, Any]
) -> None:
    geometry = _get_optional_geometry(spec)
    if not geometry:
        return

    shape = editor._get_shape_by_id(slide_id, shape_id)
    if "left" in geometry:
        shape.left = Inches(geometry["left"])
    if "top" in geometry:
        shape.top = Inches(geometry["top"])
    if "width" in geometry:
        shape.width = Inches(geometry["width"])
    if "height" in geometry:
        shape.height = Inches(geometry["height"])


def _record_named_shape(
    result: Dict[str, Any], spec: Dict[str, Any], shape_id: int
) -> None:
    name = spec.get("name")
    if name is None:
        return
    if name in result["named_shapes"]:
        raise ValueError(f"Duplicate shape name '{name}' in one tool call")
    result["named_shapes"][name] = shape_id


def _create_accent_bar(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    return editor.add_accent_bar_by_id(
        slide_id,
        _get_required(spec, "color_hex"),
        height=spec.get("height", 0.35),
        theme_name=theme_name,
        bind_theme=True,
    )


def _create_text(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_id = editor.add_textbox_by_id(
        slide_id,
        _get_required(spec, "x"),
        _get_required(spec, "y"),
        _get_required(spec, "w"),
        _get_required(spec, "h"),
    )
    editor.insert_text_by_id(slide_id, shape_id, _get_required(spec, "text"))
    if spec.get("style"):
        editor.style_text_by_id(
            slide_id,
            shape_id,
            theme_name=theme_name,
            bind_theme=True,
            **spec["style"],
        )
    return shape_id


def _create_citation(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_id = editor.add_textbox_by_id(
        slide_id,
        spec.get("x", 0.5),
        spec.get("y", 6.8),
        spec.get("w", 9.0),
        spec.get("h", 0.4),
    )
    editor.insert_text_by_id(slide_id, shape_id, _get_required(spec, "text"))

    style = {"font_size_pt": 10, "italic": True}
    style.update(spec.get("style", {}))
    editor.style_text_by_id(
        slide_id,
        shape_id,
        theme_name=theme_name,
        bind_theme=True,
        **style,
    )
    return shape_id


def _create_chart(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    return editor.add_chart_by_id(
        slide_id,
        _get_required(spec, "chart_type"),
        _get_required(spec, "chart_data"),
        _get_required(spec, "x"),
        _get_required(spec, "y"),
        _get_required(spec, "w"),
        _get_required(spec, "h"),
        style=spec.get("style"),
        theme_name=theme_name,
        bind_theme=True,
    )


def _create_table(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    table_data = _get_required(spec, "table_data")
    rows = len(table_data)
    cols = max((len(row) for row in table_data), default=0)
    if rows == 0 or cols == 0:
        raise ValueError("table_data must include at least one row and one column")

    return editor.add_table_by_id(
        slide_id,
        rows,
        cols,
        table_data,
        _get_required(spec, "x"),
        _get_required(spec, "y"),
        _get_required(spec, "w"),
        _get_required(spec, "h"),
        style=spec.get("style"),
        theme_name=theme_name,
        bind_theme=True,
    )


def _create_image(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    del theme_name
    return editor.add_image_by_id(
        slide_id,
        _get_required(spec, "image_path"),
        _get_required(spec, "x"),
        _get_required(spec, "y"),
        cx=spec.get("w"),
        cy=spec.get("h"),
    )


def _update_text(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_id = _get_required(spec, "shape_id")
    _apply_geometry_updates(editor, slide_id, shape_id, spec)
    if "text" in spec:
        editor.insert_text_by_id(slide_id, shape_id, spec["text"])
    if "style" in spec:
        _remove_bindings(
            editor,
            slide_id=slide_id,
            shape_id=shape_id,
            apply_method="style_text_by_id",
        )
        editor.style_text_by_id(
            slide_id,
            shape_id,
            theme_name=theme_name,
            bind_theme=True,
            **spec["style"],
        )
    return shape_id


def _update_citation(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_id = _update_text(editor, slide_id, spec, theme_name)
    shape = editor._get_shape_by_id(slide_id, shape_id)
    if not shape.has_text_frame:
        raise ValueError("Shape has no text frame.")
    return shape_id


def _update_chart(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_id = _get_required(spec, "shape_id")
    _apply_geometry_updates(editor, slide_id, shape_id, spec)
    if "style" in spec:
        _remove_bindings(
            editor,
            slide_id=slide_id,
            shape_id=shape_id,
            apply_method="style_chart_by_id",
        )
        editor.style_chart_by_id(
            slide_id,
            shape_id,
            theme_name=theme_name,
            bind_theme=True,
            **spec["style"],
        )
    return shape_id


def _update_table(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_id = _get_required(spec, "shape_id")
    _apply_geometry_updates(editor, slide_id, shape_id, spec)
    if "style" in spec:
        _remove_bindings(
            editor,
            slide_id=slide_id,
            shape_id=shape_id,
            apply_method="style_table_by_id",
        )
        editor.style_table_by_id(
            slide_id,
            shape_id,
            theme_name=theme_name,
            bind_theme=True,
            **spec["style"],
        )
    return shape_id


def _update_accent_bar(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_id = _get_required(spec, "shape_id")
    _apply_geometry_updates(editor, slide_id, shape_id, spec)
    if "height" in spec:
        shape = editor._get_shape_by_id(slide_id, shape_id)
        shape.height = Inches(spec["height"])
    if "color_hex" in spec:
        _remove_bindings(
            editor,
            slide_id=slide_id,
            shape_id=shape_id,
            apply_method="style_shape_fill_by_id",
        )
        editor.style_shape_fill_by_id(
            slide_id,
            shape_id,
            fill_color_hex=spec["color_hex"],
            theme_name=theme_name,
            bind_theme=True,
        )
    return shape_id


def _update_image(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    del theme_name
    shape_id = _get_required(spec, "shape_id")
    if "image_path" in spec:
        raise ValueError(
            "Updating image content is not supported; delete and re-add the image"
        )
    _apply_geometry_updates(editor, slide_id, shape_id, spec)
    return shape_id


_CREATE_DISPATCH = {
    "accent_bar": _create_accent_bar,
    "text": _create_text,
    "citation": _create_citation,
    "chart": _create_chart,
    "table": _create_table,
    "image": _create_image,
}

_UPDATE_DISPATCH = {
    "accent_bar": _update_accent_bar,
    "text": _update_text,
    "citation": _update_citation,
    "chart": _update_chart,
    "table": _update_table,
    "image": _update_image,
}


def _create_shape_from_spec(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_type = _get_required(spec, "type")
    if shape_type not in _CREATE_DISPATCH:
        raise ValueError(f"Unsupported shape type '{shape_type}'")
    return _CREATE_DISPATCH[shape_type](editor, slide_id, spec, theme_name)


def _update_shape_from_spec(
    editor: PptxEditor, slide_id: int, spec: Dict[str, Any], theme_name: str
) -> int:
    shape_type = _get_required(spec, "type")
    if shape_type not in _UPDATE_DISPATCH:
        raise ValueError(f"Unsupported shape type '{shape_type}'")
    return _UPDATE_DISPATCH[shape_type](editor, slide_id, spec, theme_name)


def _delete_shape(editor: PptxEditor, slide_id: int, shape_id: int) -> None:
    shape = editor._get_shape_by_id(slide_id, shape_id)
    shape._element.getparent().remove(shape._element)
    _remove_bindings(editor, slide_id=slide_id, shape_id=shape_id)


def create_slide(
    editor: PptxEditor,
    *,
    layout_index: int = 6,
    theme_name: str = "default",
    background_color: Optional[str] = None,
    shapes: Optional[list[Dict[str, Any]]] = None,
) -> Dict[str, Any]:
    slide_index = editor.add_slide(layout_index)
    slide_id = editor.get_slide_id(slide_index)
    result = {"slide_id": slide_id, "shape_ids": [], "named_shapes": {}}

    if background_color is not None:
        editor.set_slide_background_by_id(
            slide_id,
            background_color,
            theme_name=theme_name,
            bind_theme=True,
        )

    for spec in shapes or []:
        shape_id = _create_shape_from_spec(editor, slide_id, spec, theme_name)
        result["shape_ids"].append(shape_id)
        _record_named_shape(result, spec, shape_id)

    return result


def update_slide(
    editor: PptxEditor,
    slide_id: int,
    *,
    theme_name: str = "default",
    background_color: Optional[str] = None,
    delete_shape_ids: Optional[list[int]] = None,
    add_shapes: Optional[list[Dict[str, Any]]] = None,
    update_shapes: Optional[list[Dict[str, Any]]] = None,
) -> Dict[str, Any]:
    editor._get_slide_by_id(slide_id)
    delete_shape_ids = delete_shape_ids or []
    add_shapes = add_shapes or []
    update_shapes = update_shapes or []

    delete_set = set(delete_shape_ids)
    if len(delete_set) != len(delete_shape_ids):
        raise ValueError("delete_shape_ids contains duplicates")

    update_shape_ids = []
    for spec in update_shapes:
        shape_id = _get_required(spec, "shape_id")
        if shape_id in delete_set:
            raise ValueError(
                f"Shape id {shape_id} cannot be both updated and deleted in one call"
            )
        update_shape_ids.append(shape_id)

    if len(set(update_shape_ids)) != len(update_shape_ids):
        raise ValueError("update_shapes contains duplicate shape_id values")

    result = {
        "slide_id": slide_id,
        "deleted_shape_ids": [],
        "updated_shape_ids": [],
        "created_shape_ids": [],
        "named_shapes": {},
    }

    if background_color is not None:
        _remove_bindings(
            editor,
            slide_id=slide_id,
            shape_id=None,
            apply_method="set_slide_background_by_id",
        )
        editor.set_slide_background_by_id(
            slide_id,
            background_color,
            theme_name=theme_name,
            bind_theme=True,
        )

    for spec in update_shapes:
        shape_id = _update_shape_from_spec(editor, slide_id, spec, theme_name)
        result["updated_shape_ids"].append(shape_id)
        _record_named_shape(result, spec, shape_id)

    for spec in add_shapes:
        shape_id = _create_shape_from_spec(editor, slide_id, spec, theme_name)
        result["created_shape_ids"].append(shape_id)
        _record_named_shape(result, spec, shape_id)

    for shape_id in delete_shape_ids:
        _delete_shape(editor, slide_id, shape_id)
        result["deleted_shape_ids"].append(shape_id)

    return result


__all__ = [
    "register_theme",
    "update_theme",
    "create_slide",
    "update_slide",
]
