from pathlib import Path
from typing import Any, Dict, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from src.utils.pptx_functions import PptxEditor


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def create_presentation(path: Optional[str] = None) -> PptxEditor:
    return PptxEditor(path)


def save_presentation(editor: PptxEditor, output_path: str) -> str:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    editor.prs.save(output)
    return str(output)


def get_theme_preset(name: str) -> Dict[str, str]:
    presets = {
        "modern_business": {
            "bg": "#F6F9FC",
            "surface": "#FFFFFF",
            "accent": "#0F62FE",
            "primary": "#102A43",
            "secondary": "#486581",
            "font": "Aptos",
        },
        "dark_analytics": {
            "bg": "#0B1020",
            "surface": "#131A2A",
            "accent": "#4FD1C5",
            "primary": "#F7FAFC",
            "secondary": "#A0AEC0",
            "font": "Aptos Display",
        },
        "academic_report": {
            "bg": "#F8F5EF",
            "surface": "#FFFDF9",
            "accent": "#7A5C3E",
            "primary": "#2D241D",
            "secondary": "#6B5B4D",
            "font": "Georgia",
        },
    }
    if name not in presets:
        raise ValueError(f"Unknown theme preset: {name}")
    return dict(presets[name])


def apply_slide_background(
    editor: PptxEditor, slide_index: int, color_hex: str
) -> None:
    editor._validate_slide(slide_index)
    slide = editor.prs.slides[slide_index]
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)


def add_accent_bar(
    editor: PptxEditor,
    slide_index: int,
    color_hex: str,
    *,
    height: float = 0.35,
) -> int:
    editor._validate_slide(slide_index)
    slide = editor.prs.slides[slide_index]
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(height),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(color_hex)
    accent.line.fill.background()
    return len(slide.shapes) - 1


def apply_slide_theme(
    editor: PptxEditor,
    slide_index: int,
    *,
    background_color: Optional[str] = None,
    accent_color: Optional[str] = None,
) -> None:
    if background_color is not None:
        apply_slide_background(editor, slide_index, background_color)
    if accent_color is not None:
        add_accent_bar(editor, slide_index, accent_color)


def add_blank_slide(
    editor: PptxEditor,
    *,
    background_color: Optional[str] = None,
    accent_color: Optional[str] = None,
    layout_index: int = 6,
) -> int:
    slide_index = editor.add_slide(layout_index)
    apply_slide_theme(
        editor,
        slide_index,
        background_color=background_color,
        accent_color=accent_color,
    )
    return slide_index


def add_textbox(
    editor: PptxEditor, slide_index: int, x: float, y: float, w: float, h: float
) -> int:
    editor._validate_slide(slide_index)
    slide = editor.prs.slides[slide_index]
    slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return len(slide.shapes) - 1


def add_text_block(
    editor: PptxEditor,
    slide_index: int,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_name: str,
    font_size: int,
    color_hex: str,
    bold: bool = False,
    italic: bool = False,
) -> int:
    shape_index = add_textbox(editor, slide_index, x, y, w, h)
    editor.insert_text(slide_index, shape_index, text)
    editor.style_update(
        slide_index,
        shape_index,
        font_name=font_name,
        font_size_pt=font_size,
        bold=bold,
        italic=italic,
        color_hex=color_hex,
    )
    return shape_index


def add_chart_block(
    editor: PptxEditor,
    slide_index: int,
    chart_type: str,
    chart_data: Dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    style: Optional[Dict[str, Any]] = None,
) -> int:
    editor._validate_slide(slide_index)
    slide = editor.prs.slides[slide_index]
    editor.add_chart(slide_index, chart_type, chart_data, x, y, w, h, style=style)
    return len(slide.shapes) - 1


def add_table_block(
    editor: PptxEditor,
    slide_index: int,
    table_data: list[list[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    style: Optional[Dict[str, Any]] = None,
) -> int:
    rows = len(table_data)
    cols = max((len(row) for row in table_data), default=0)
    if rows == 0 or cols == 0:
        raise ValueError("table_data must include at least one row and one column")

    editor._validate_slide(slide_index)
    slide = editor.prs.slides[slide_index]
    editor.add_table(slide_index, rows, cols, table_data, x, y, w, h, style=style)
    return len(slide.shapes) - 1


def add_citation_block(
    editor: PptxEditor,
    slide_index: int,
    text: str,
    *,
    font_name: Optional[str] = None,
    font_size_pt: Optional[int] = None,
    italic: Optional[bool] = True,
    color_hex: Optional[str] = None,
) -> int:
    editor._validate_slide(slide_index)
    slide = editor.prs.slides[slide_index]
    editor.add_citation(slide_index, text)
    shape_index = len(slide.shapes) - 1
    if any(value is not None for value in (font_name, font_size_pt, italic, color_hex)):
        editor.style_update(
            slide_index,
            shape_index,
            font_name=font_name,
            font_size_pt=font_size_pt,
            italic=italic,
            color_hex=color_hex,
        )
    return shape_index


__all__ = [
    "PptxEditor",
    "rgb",
    "create_presentation",
    "save_presentation",
    "get_theme_preset",
    "apply_slide_background",
    "add_accent_bar",
    "apply_slide_theme",
    "add_blank_slide",
    "add_textbox",
    "add_text_block",
    "add_chart_block",
    "add_table_block",
    "add_citation_block",
]
