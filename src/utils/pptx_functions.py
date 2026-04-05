import os
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


# a bit of own, a bit of vibe, makes everyone happy
class PptxEditor:
    def __init__(self, path: Optional[str] = None):
        self.prs = Presentation(path) if path else Presentation()

    def _validate_slide(self, index: int):
        if index < 0 or index >= len(self.prs.slides):
            raise IndexError(f"Slide index {index} out of range")

    def _get_shape(self, slide_index: int, shape_index: int):
        self._validate_slide(slide_index)
        shapes = list(self.prs.slides[slide_index].shapes)
        if shape_index < 0 or shape_index >= len(shapes):
            raise IndexError(f"Shape index {shape_index} out of range")
        return shapes[shape_index]

    def _rgb_color(self, color_hex: str) -> RGBColor:
        return RGBColor.from_string(color_hex.replace("#", ""))

    def _apply_font_style(
        self,
        font,
        font_name: Optional[str] = None,
        font_size_pt: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color_hex: Optional[str] = None,
    ):
        if font_name is not None:
            font.name = font_name
        if font_size_pt is not None:
            font.size = Pt(font_size_pt)
        if bold is not None:
            font.bold = bold
        if italic is not None:
            font.italic = italic
        if color_hex is not None:
            font.color.rgb = self._rgb_color(color_hex)

    def _apply_text_frame_style(
        self,
        text_frame,
        font_name: Optional[str] = None,
        font_size_pt: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color_hex: Optional[str] = None,
    ):
        for paragraph in text_frame.paragraphs:
            self._apply_font_style(
                paragraph.font,
                font_name=font_name,
                font_size_pt=font_size_pt,
                bold=bold,
                italic=italic,
                color_hex=color_hex,
            )
            for run in paragraph.runs:
                self._apply_font_style(
                    run.font,
                    font_name=font_name,
                    font_size_pt=font_size_pt,
                    bold=bold,
                    italic=italic,
                    color_hex=color_hex,
                )

    def add_slide(self, layout_index: int = 1) -> int:
        layout_index = max(0, min(layout_index, len(self.prs.slide_layouts) - 1))
        layout = self.prs.slide_layouts[layout_index]
        self.prs.slides.add_slide(layout)
        return len(self.prs.slides) - 1

    def reorder_slide(self, old_index: int, new_index: int):
        self._validate_slide(old_index)
        xml_slides = self.prs.slides._sldIdLst
        if new_index < 0 or new_index >= len(xml_slides):
            raise IndexError("new_index out of bounds")
        slide_elem = xml_slides[old_index]
        del xml_slides[old_index]
        xml_slides.insert(new_index, slide_elem)

    def set_slide_layout(self, slide_index: int, layout_index: int):
        self._validate_slide(slide_index)
        if layout_index < 0 or layout_index >= len(self.prs.slide_layouts):
            raise IndexError("Invalid layout index")
        slide = self.prs.slides[slide_index]
        new_layout = self.prs.slide_layouts[layout_index]

        layout_rel = None
        for rel in slide.part.rels.values():
            if "slideLayout" in rel.reltype:
                layout_rel = rel
                break

        if layout_rel:
            slide.part.drop_rel(layout_rel.rId)

        slide.part.relate_to(
            new_layout.part,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        )
        if hasattr(slide, "_slide_layout"):
            del slide._slide_layout

    def insert_text(self, slide_index: int, shape_index: int, text: str):
        shape = self._get_shape(slide_index, shape_index)
        if not shape.has_text_frame:
            raise ValueError("Shape has no text frame.")
        shape.text_frame.text = str(text)

    def add_chart(
        self,
        slide_index: int,
        chart_type_str: str,
        chart_data_dict: Dict[str, Any],
        x: float,
        y: float,
        cx: float,
        cy: float,
        style: Optional[Dict[str, Any]] = None,
    ):
        self._validate_slide(slide_index)
        slide = self.prs.slides[slide_index]

        chart_data = CategoryChartData()
        chart_data.categories = chart_data_dict.get("categories", [])
        for series in chart_data_dict.get("series", []):
            chart_data.add_series(series["name"], tuple(series["values"]))

        chart_type_enum = getattr(
            XL_CHART_TYPE, chart_type_str.upper(), XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        slide.shapes.add_chart(
            chart_type_enum, Inches(x), Inches(y), Inches(cx), Inches(cy), chart_data
        )
        if style:
            self.style_chart(slide_index, len(slide.shapes) - 1, **style)

    def add_table(
        self,
        slide_index: int,
        rows: int,
        cols: int,
        table_data: List[List[str]],
        x: float,
        y: float,
        cx: float,
        cy: float,
        style: Optional[Dict[str, Any]] = None,
    ):
        self._validate_slide(slide_index)
        slide = self.prs.slides[slide_index]
        table = slide.shapes.add_table(
            rows, cols, Inches(x), Inches(y), Inches(cx), Inches(cy)
        ).table

        for r in range(min(rows, len(table_data))):
            for c in range(min(cols, len(table_data[r]))):
                table.cell(r, c).text = str(table_data[r][c])
        if style:
            self.style_table(slide_index, len(slide.shapes) - 1, **style)

    def add_image(
        self,
        slide_index: int,
        image_path: str,
        x: float,
        y: float,
        cx: Optional[float] = None,
        cy: Optional[float] = None,
    ):
        self._validate_slide(slide_index)
        if not os.path.isfile(image_path):
            raise FileNotFoundError("Image path invalid.")

        slide = self.prs.slides[slide_index]
        kwargs = {"left": Inches(x), "top": Inches(y)}
        if cx is not None:
            kwargs["width"] = Inches(cx)
        if cy is not None:
            kwargs["height"] = Inches(cy)

        slide.shapes.add_picture(image_path, **kwargs)

    def add_citation(self, slide_index: int, text: str):
        self._validate_slide(slide_index)
        slide = self.prs.slides[slide_index]
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.4)
        )
        p = txBox.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.italic = True

    def style_update(
        self,
        slide_index: int,
        shape_index: int,
        font_name: Optional[str] = None,
        font_size_pt: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color_hex: Optional[str] = None,
    ):
        shape = self._get_shape(slide_index, shape_index)
        if not shape.has_text_frame:
            raise ValueError("Shape has no text frame.")

        self._apply_text_frame_style(
            shape.text_frame,
            font_name=font_name,
            font_size_pt=font_size_pt,
            bold=bold,
            italic=italic,
            color_hex=color_hex,
        )

    def style_table(
        self,
        slide_index: int,
        shape_index: int,
        header_fill_hex: Optional[str] = None,
        body_fill_hex: Optional[str] = None,
        header_font_name: Optional[str] = None,
        body_font_name: Optional[str] = None,
        header_font_size_pt: Optional[int] = None,
        body_font_size_pt: Optional[int] = None,
        header_bold: Optional[bool] = True,
        body_bold: Optional[bool] = None,
        header_italic: Optional[bool] = None,
        body_italic: Optional[bool] = None,
        header_font_color_hex: Optional[str] = None,
        body_font_color_hex: Optional[str] = None,
    ):
        shape = self._get_shape(slide_index, shape_index)
        if not shape.has_table:
            raise ValueError("Shape has no table.")

        table = shape.table
        for row_idx, row in enumerate(table.rows):
            is_header = row_idx == 0
            fill_hex = header_fill_hex if is_header else body_fill_hex
            font_name = header_font_name if is_header else body_font_name
            font_size_pt = header_font_size_pt if is_header else body_font_size_pt
            bold = header_bold if is_header else body_bold
            italic = header_italic if is_header else body_italic
            font_color_hex = header_font_color_hex if is_header else body_font_color_hex

            for cell in row.cells:
                if fill_hex is not None:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._rgb_color(fill_hex)
                self._apply_text_frame_style(
                    cell.text_frame,
                    font_name=font_name,
                    font_size_pt=font_size_pt,
                    bold=bold,
                    italic=italic,
                    color_hex=font_color_hex,
                )

    def style_chart(
        self,
        slide_index: int,
        shape_index: int,
        title: Optional[str] = None,
        title_font_name: Optional[str] = None,
        title_font_size_pt: Optional[int] = None,
        title_bold: Optional[bool] = None,
        title_italic: Optional[bool] = None,
        title_color_hex: Optional[str] = None,
        legend_font_name: Optional[str] = None,
        legend_font_size_pt: Optional[int] = None,
        legend_bold: Optional[bool] = None,
        legend_italic: Optional[bool] = None,
        legend_color_hex: Optional[str] = None,
        axis_font_name: Optional[str] = None,
        axis_font_size_pt: Optional[int] = None,
        axis_bold: Optional[bool] = None,
        axis_italic: Optional[bool] = None,
        axis_color_hex: Optional[str] = None,
        series_colors: Optional[List[str]] = None,
    ):
        shape = self._get_shape(slide_index, shape_index)
        if not shape.has_chart:
            raise ValueError("Shape has no chart.")

        chart = shape.chart

        if title is not None:
            chart.has_title = True
            chart.chart_title.text_frame.text = str(title)
        if chart.has_title:
            self._apply_text_frame_style(
                chart.chart_title.text_frame,
                font_name=title_font_name,
                font_size_pt=title_font_size_pt,
                bold=title_bold,
                italic=title_italic,
                color_hex=title_color_hex,
            )

        if chart.has_legend:
            self._apply_font_style(
                chart.legend.font,
                font_name=legend_font_name,
                font_size_pt=legend_font_size_pt,
                bold=legend_bold,
                italic=legend_italic,
                color_hex=legend_color_hex,
            )

        for axis_name in ("category_axis", "value_axis"):
            axis = getattr(chart, axis_name, None)
            if axis is not None and hasattr(axis, "tick_labels"):
                self._apply_font_style(
                    axis.tick_labels.font,
                    font_name=axis_font_name,
                    font_size_pt=axis_font_size_pt,
                    bold=axis_bold,
                    italic=axis_italic,
                    color_hex=axis_color_hex,
                )

        if series_colors:
            for series, color_hex in zip(chart.series, series_colors):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = self._rgb_color(color_hex)
                series.format.line.color.rgb = self._rgb_color(color_hex)


if __name__ == "__main__":
    pptx_editor = PptxEditor()
    pptx_editor.add_slide()
    pptx_editor.insert_text(0, 0, "Hello, World!")
    pptx_editor.insert_text(0, 1, "Item 1\nItem 2\nItem 3")

    ppt = pptx_editor.prs
    os.makedirs("outputs", exist_ok=True)
    ppt.save("outputs/output.pptx")
