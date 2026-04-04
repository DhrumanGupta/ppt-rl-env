import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN
from typing import Any, Dict, List, Optional

# a bit of own, a bit of vibe, makes everyone happy
class PptxEditor:

    def __init__(self, path: Optional[str] = None):
        self.prs = Presentation(path) if path else Presentation()

    def _validate_slide(self, index: int):
        if index < 0 or index >= len(self.prs.slides):
            raise IndexError(f"Slide index {index} out of range")

    def _get_shape(self, slide_index: int, shape_index: int):
        self._validate_slide(slide_index)
        shapes = list(self.prs.slides[slide_index].shapes)
        if shape_index < 0 or shape_index >= len(shapes):
            raise IndexError(f"Shape index {shape_index} out of range")
        return shapes[shape_index]

    def add_slide(self, layout_index: int = 1) -> int:
        layout_index = max(0, min(layout_index, len(self.prs.slide_layouts) - 1))
        layout = self.prs.slide_layouts[layout_index]
        self.prs.slides.add_slide(layout)
        return len(self.prs.slides) - 1

    def reorder_slide(self, old_index: int, new_index: int):
        self._validate_slide(old_index)
        xml_slides = self.prs.slides._sldIdLst
        if new_index < 0 or new_index >= len(xml_slides):
            raise IndexError("new_index out of bounds")
        slide_elem = xml_slides[old_index]
        del xml_slides[old_index]
        xml_slides.insert(new_index, slide_elem)

    def set_slide_layout(self, slide_index: int, layout_index: int):
        self._validate_slide(slide_index)
        if layout_index < 0 or layout_index >= len(self.prs.slide_layouts):
            raise IndexError("Invalid layout index")
        slide = self.prs.slides[slide_index]
        new_layout = self.prs.slide_layouts[layout_index]

        layout_rel = None
        for rel in slide.part.rels.values():
            if "slideLayout" in rel.reltype:
                layout_rel = rel
                break

        if layout_rel:
            slide.part.drop_rel(layout_rel.rId)

        slide.part.relate_to(
            new_layout.part,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
        )
        if hasattr(slide, '_slide_layout'):
            del slide._slide_layout

    def write_title(self, slide_index: int, text: str):
        self._validate_slide(slide_index)
        slide = self.prs.slides[slide_index]

        if slide.shapes.title:
            slide.shapes.title.text = text
        else:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            tf = txBox.text_frame
            tf.text = text
            tf.paragraphs[0].font.size = Pt(36)
            tf.paragraphs[0].font.bold = True

    def write_bullets(self, slide_index: int, shape_index: int, bullets: List[str]):
        shape = self._get_shape(slide_index, shape_index)
        if not shape.has_text_frame:
            raise ValueError("Target shape does not support text.")

        tf = shape.text_frame
        tf.clear()

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(bullet)
            p.level = 0

    def add_chart(self, slide_index: int, chart_type_str: str, chart_data_dict: Dict[str, Any],
                  x: float, y: float, cx: float, cy: float):
        self._validate_slide(slide_index)
        slide = self.prs.slides[slide_index]

        chart_data = CategoryChartData()
        chart_data.categories = chart_data_dict.get("categories", [])
        for series in chart_data_dict.get("series", []):
            chart_data.add_series(series["name"], tuple(series["values"]))

        chart_type_enum = getattr(XL_CHART_TYPE, chart_type_str.upper(), XL_CHART_TYPE.COLUMN_CLUSTERED)
        slide.shapes.add_chart(chart_type_enum, Inches(x), Inches(y), Inches(cx), Inches(cy), chart_data)

    def add_table(self, slide_index: int, rows: int, cols: int, table_data: List[List[str]],
                  x: float, y: float, cx: float, cy: float):
        self._validate_slide(slide_index)
        slide = self.prs.slides[slide_index]
        table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(cx), Inches(cy)).table

        for r in range(min(rows, len(table_data))):
            for c in range(min(cols, len(table_data[r]))):
                table.cell(r, c).text = str(table_data[r][c])

    def add_image(self, slide_index: int, image_path: str, x: float, y: float,
                  cx: Optional[float] = None, cy: Optional[float] = None):
        self._validate_slide(slide_index)
        if not os.path.isfile(image_path):
            raise FileNotFoundError("Image path invalid.")

        slide = self.prs.slides[slide_index]
        kwargs = {'left': Inches(x), 'top': Inches(y)}
        if cx is not None: kwargs['width'] = Inches(cx)
        if cy is not None: kwargs['height'] = Inches(cy)

        slide.shapes.add_picture(image_path, **kwargs)

    def add_citation(self, slide_index: int, text: str):
        self._validate_slide(slide_index)
        slide = self.prs.slides[slide_index]
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.italic = True

    def style_update(self, slide_index: int, shape_index: int,
                     font_name: Optional[str] = None, font_size_pt: Optional[int] = None,
                     bold: Optional[bool] = None, italic: Optional[bool] = None,
                     color_hex: Optional[str] = None):
        shape = self._get_shape(slide_index, shape_index)
        if not shape.has_text_frame:
            raise ValueError("Shape has no text frame.")

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if font_name is not None: run.font.name = font_name
                if font_size_pt is not None: run.font.size = Pt(font_size_pt)
                if bold is not None: run.font.bold = bold
                if italic is not None: run.font.italic = italic
                if color_hex is not None: run.font.color.rgb = RGBColor.from_string(color_hex.replace("#", ""))

    def refine_text(self, slide_index: int, shape_index: int, new_text: str):
        shape = self._get_shape(slide_index, shape_index)
        if not shape.has_text_frame:
            raise ValueError("Shape has no text frame.")
        shape.text_frame.text = new_text
